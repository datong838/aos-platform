#!/usr/bin/env python3
"""
Full PPTX recreation: Convert all 37 slides of 讲透AOS企业AI操作系统.pptx
from images to native PPTX components with white background.

All content is reconstructed using shapes (rounded rectangles, ovals, arrows,
text boxes, connectors) instead of embedded images.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = "/Users/ddt/work/projects/ai_agent/docs/ref/讲透AOS企业AI操作系统-白色版.pptx"

# ============ Color Palette ============
COL_TITLE_DARK = RGBColor(0x10, 0x35, 0x5E)
COL_TITLE_NAVY = RGBColor(0x0D, 0x21, 0x3E)
COL_SUB_BLUE = RGBColor(0x1F, 0x6F, 0xC1)
COL_BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_LIGHT_BLUE = RGBColor(0xE9, 0xF2, 0xFC)
COL_LIGHTER_BLUE = RGBColor(0xF0, 0xF6, 0xFE)
COL_SECTION_HEADER = RGBColor(0x12, 0x47, 0x9A)
COL_ICON_BLUE = RGBColor(0x1F, 0x6F, 0xC1)
COL_BORDER_BLUE = RGBColor(0xC5, 0xDB, 0xF0)
COL_MEDIUM_BLUE = RGBColor(0x4D, 0x88, 0xD6)
COL_TEXT_DARK = RGBColor(0x1F, 0x2D, 0x3D)
COL_TEXT_GRAY = RGBColor(0x5A, 0x6B, 0x7E)
COL_TEXT_LIGHT = RGBColor(0x8A, 0x95, 0xA5)
COL_GREEN = RGBColor(0x2D, 0xC2, 0x76)
COL_GREEN_BG = RGBColor(0xE8, 0xF8, 0xF0)
COL_RED = RGBColor(0xE5, 0x4A, 0x4A)
COL_RED_BG = RGBColor(0xFD, 0xEC, 0xEC)
COL_ORANGE = RGBColor(0xF5, 0x8A, 0x1F)
COL_ORANGE_BG = RGBColor(0xFE, 0xF3, 0xE5)
COL_PURPLE = RGBColor(0x7C, 0x5B, 0xD8)
COL_PURPLE_BG = RGBColor(0xF3, 0xEF, 0xFC)
COL_TEAL = RGBColor(0x0D, 0x9B, 0x8A)
COL_TEAL_BG = RGBColor(0xE5, 0xF7, 0xF5)
COL_LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
COL_BORDER_GRAY = RGBColor(0xD8, 0xDF, 0xE8)
COL_BORDER_LIGHT = RGBColor(0xEC, 0xF0, 0xF5)
COL_ARROW_GRAY = RGBColor(0xB0, 0xBD, 0xCC)
COL_VALUE_BG = RGBColor(0xFA, 0xFB, 0xFD)
COL_GOLD = RGBColor(0xC4, 0x9A, 0x2A)
COL_DARK_BG = RGBColor(0x0D, 0x21, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ============ Helper Functions ============

def set_bg_white(s):
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COL_BG_WHITE

def add_text(slide, x, y, w, h, text, size=12, color=None, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Microsoft YaHei",
             line_spacing=1.15):
    if color is None:
        color = COL_TEXT_DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(text, list):
        for i, line in enumerate(text):
            if i > 0:
                p = tf.add_paragraph()
                p.alignment = align
                if line_spacing:
                    p.line_spacing = line_spacing
            if isinstance(line, tuple):
                run = p.add_run()
                run.text = line[0]
                run.font.size = Pt(line[1] if len(line) > 1 else size)
                run.font.color.rgb = line[2] if len(line) > 2 else color
                run.font.bold = line[3] if len(line) > 3 else bold
                run.font.name = font_name
            elif isinstance(line, str):
                run = p.add_run()
                run.text = line
                run.font.size = Pt(size)
                run.font.color.rgb = color
                run.font.bold = bold
                run.font.name = font_name
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return tb

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5,
             radius=0.25):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    sp = shape._element
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        radius_val = int(radius * 100000)
        gd.set('fmla', f'val {radius_val}')
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_plain_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_arrow(slide, x, y, w, h, color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        color = COL_ARROW_GRAY
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_down_arrow(slide, x, y, w, h, color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        color = COL_ARROW_GRAY
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, x, y, d, fill_color=None, line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x), Inches(y), Inches(d), Inches(d))
    shape.shadow.inherit = False
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=None, width=1.0):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    if color is None:
        color = COL_BORDER_BLUE
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn

def add_page_number(slide, num, total="17"):
    add_text(slide, 11.8, 7.05, 1.3, 0.3,
             f"{num:02d} / {total}", size=9, color=COL_TEXT_LIGHT, align=PP_ALIGN.RIGHT)

def add_footer(slide):
    add_text(slide, 0.3, 7.15, 5, 0.25,
             "Palantir 深度解析", size=8, color=COL_TEXT_LIGHT)

def add_card(slide, x, y, w, h, title, items, accent_color, bg_color=None):
    """Add a card with colored header and bullet items."""
    if bg_color is None:
        bg_color = COL_BG_WHITE
    card = add_rect(slide, x, y, w, h, fill_color=bg_color, line_color=COL_BORDER_GRAY, line_width=0.5)
    # Accent bar at top
    add_plain_rect(slide, x, y, w, 0.06, fill_color=accent_color)
    # Title
    add_text(slide, x + 0.1, y + 0.12, w - 0.2, 0.35,
             title, size=12, color=accent_color, bold=True)
    # Items
    item_y = y + 0.5
    for item in items:
        add_circle(slide, x + 0.15, item_y + 0.06, 0.06, fill_color=accent_color)
        add_text(slide, x + 0.3, item_y, w - 0.4, 0.25,
                 item, size=9, color=COL_TEXT_GRAY)
        item_y += 0.3
    return card

def add_metric_box(slide, x, y, w, h, value, label, val_color=None):
    if val_color is None:
        val_color = COL_ICON_BLUE
    add_rect(slide, x, y, w, h, fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
    add_text(slide, x + 0.05, y + 0.05, w - 0.1, h * 0.5,
             value, size=18, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.05, y + h * 0.55, w - 0.1, h * 0.4,
             label, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

def add_title_bar(slide, num, total, title_text, subtitle=None):
    """Standard slide header with page number and title."""
    add_text(slide, 0.3, 0.15, 1.0, 0.3,
             f"{num:02d} / {total}", size=9, color=COL_TEXT_LIGHT)
    add_text(slide, 0.3, 0.4, 11.5, 0.6,
             title_text, size=24, color=COL_TITLE_DARK, bold=True)
    if subtitle:
        add_text(slide, 0.3, 1.0, 11.5, 0.35,
                 subtitle, size=12, color=COL_SUB_BLUE)
    add_line(slide, 0.3, 1.4, 13.0, 1.4, color=COL_BORDER_BLUE, width=1.0)
    add_page_number(slide, num, total)
    add_footer(slide)


# ====================================================================
# SLIDE 1: COVER
# ====================================================================
def slide_01_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    # Decorative top bar
    add_plain_rect(s, 0, 0, 13.33, 0.08, fill_color=COL_SECTION_HEADER)
    # Main title
    add_text(s, 1.0, 1.5, 11.3, 1.0,
             "23页PPT讲透 Palantir", size=40, color=COL_TITLE_DARK, bold=True,
             align=PP_ALIGN.CENTER)
    # Subtitle
    add_text(s, 1.0, 2.6, 11.3, 0.5,
             "Ontology · Data · Logic · Action · Security",
             size=18, color=COL_SUB_BLUE, align=PP_ALIGN.CENTER)
    # Description
    add_text(s, 2.0, 3.4, 9.3, 1.5,
             "以 Ontology 为业务运行模型，以 Foundry 为数据与运营平台，\n"
             "以 AIP 为 AI 运行环境，以 Apollo 和 Rubix 为软件交付底座，\n"
             "以 FDE 为产品反馈机制，连接数据、决策和现实行动的\n"
             "企业级决策与运营基础设施。",
             size=14, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.5)
    # Bottom info
    add_text(s, 1.0, 6.2, 11.3, 0.35,
             "2026.07 · 基于公开资料整理",
             size=11, color=COL_TEXT_LIGHT, align=PP_ALIGN.CENTER)
    # Decorative bottom bar
    add_plain_rect(s, 0, 7.42, 13.33, 0.08, fill_color=COL_BORDER_BLUE)

# ====================================================================
# SLIDE 2: COVER INFOGRAPHIC (same as test slide)
# ====================================================================
def slide_02_infographic():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)

    # Title
    add_text(s, 0.3, 0.2, 12.7, 0.6,
             "Palantir 究竟是一家什么公司", size=28, color=COL_TITLE_DARK, bold=True)
    add_text(s, 0.3, 0.85, 12.7, 0.4,
             "以数据、业务逻辑和AI为核心，连接决策与现实行动的企业级决策与运营基础设施",
             size=13, color=COL_SUB_BLUE)
    add_line(s, 0.3, 1.35, 13.0, 1.35, color=COL_BORDER_BLUE, width=1.0)

    # LEFT: Intro paragraph
    add_rect(s, 0.3, 1.5, 7.4, 1.55, fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY)
    add_text(s, 0.5, 1.6, 7.2, 1.35,
             "Palantir 帮助政府与大型企业将分散的数据、复杂的业务关系、规则、模型和人员行动连接在一起，"
             "在同一套系统中完成从「感知 → 决策 → 行动 → 结果反馈」的闭环，"
             "让组织在复杂和高风险的环境中更快、更准、更安全地行动。",
             size=11, color=COL_TEXT_GRAY)

    # RIGHT: Core values
    add_text(s, 7.9, 1.4, 5.2, 0.35,
             "为客户创造的核心价值", size=13, color=COL_SECTION_HEADER, bold=True)
    values = [
        ("打破数据孤岛", "形成单一态势"),
        ("更快更准的决策", "发现、模型 + AI"),
        ("让决策落地为行动", "产生真实影响"),
        ("企业级安全与治理", "可追溯、可审计"),
        ("持续学习与优化", "形成决策资产"),
        ("跨环境部署", "云·私有·边缘"),
    ]
    for i, (title, sub) in enumerate(values):
        col = i % 3
        row = i // 3
        x = 7.9 + col * 1.78
        y = 1.8 + row * 0.7
        add_rect(s, x, y, 1.7, 0.62, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
        add_text(s, x + 0.05, y + 0.05, 1.6, 0.3, title, size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.32, 1.6, 0.3, sub, size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # ROW 2 LEFT: Capabilities
    add_text(s, 0.3, 3.25, 7.4, 0.35, "Palantir 的真实身份", size=14, color=COL_SECTION_HEADER, bold=True)
    caps = [
        ("数据库", "数据连接与治理", "连接各数据源、清洗、可追溯"),
        ("业务建模", "Ontology 表达", "对象与关系刻画"),
        ("逻辑、智能", "决策能力建模", "规则 + 模型 + AI"),
        ("Action", "行动执行", "决策改变现实"),
        ("安全权限", "安全与权限", "贯穿运营全过程"),
        ("反馈", "结果反馈", "持续学习与优化"),
    ]
    for i, (icon, title, sub) in enumerate(caps):
        x = 0.3 + i * 1.23
        y = 3.7
        add_rect(s, x, y, 1.13, 1.1, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
        add_circle(s, x + 0.42, y + 0.1, 0.3, fill_color=COL_LIGHT_BLUE)
        add_text(s, x + 0.42, y + 0.14, 0.3, 0.25, icon[:2] if len(icon) >= 2 else icon,
                 size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.02, y + 0.5, 1.09, 0.3, title, size=9, color=COL_TITLE_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.02, y + 0.78, 1.09, 0.28, sub, size=7, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # RIGHT: Key data
    add_text(s, 7.9, 3.25, 5.2, 0.35, "Palantir 关键数据（截至 2026 Q1）", size=14, color=COL_SECTION_HEADER, bold=True)
    metrics = [
        ("4,429", "全职员工", COL_ICON_BLUE),
        ("954+", "家客户", COL_GREEN),
        ("$44.75亿", "2025 年度收入", COL_RED),
        ("82%", "2025 毛利率", COL_SECTION_HEADER),
        ("85%", "2025 年毛利率", COL_MEDIUM_BLUE),
        ("87%", "2026 Q1 毛利率", COL_TITLE_DARK),
    ]
    for i, (val, label, c) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = 7.9 + col * 1.78
        y = 3.7 + row * 0.7
        add_rect(s, x, y, 1.7, 0.62, fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_text(s, x + 0.05, y + 0.05, 1.6, 0.32, val, size=16, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.4, 1.6, 0.22, label, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, 7.9, 5.18, 5.2, 0.25, "数据来源：Palantir 2025 年度报告、2026 Q1 财报", size=8, color=COL_TEXT_GRAY)

    # ROW 3: Flow
    add_text(s, 0.3, 5.55, 12.7, 0.35, "Palantir 连接什么 → 产生什么", size=14, color=COL_SECTION_HEADER, bold=True)
    steps = [
        ("连接所有数据", "内部系统(ERP/MES/CRM)\n数据库/文档/传感器/IoT\n外部数据/合作伙伴", COL_SECTION_HEADER),
        ("表达业务\n现实世界", "客户\n工厂\n产品\n供应商", COL_ICON_BLUE),
        ("决策与智能", "业务规则\n预测模型\n优化算法", COL_GREEN),
        ("行动与执行", "审批与流程\n调用外部\n调整计划", COL_SECTION_HEADER),
        ("结果反馈\n持续优化", "执行结果追踪\n效果分析\n沉淀资产", COL_RED),
    ]
    box_w = 2.4
    gap = 0.1
    total_w = box_w * 5 + gap * 4
    start_x = (13.33 - total_w) / 2
    for i, (title, content, c) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        y = 6.0
        add_rect(s, x, y, box_w, 1.05, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=1.0)
        add_text(s, x + 0.05, y + 0.05, box_w - 0.1, 0.3, title, size=11, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.32, box_w - 0.1, 0.7, content, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(s, x + box_w - 0.05, y + 0.4, gap + 0.1, 0.25)

    # Bottom summary
    add_rect(s, 0.3, 7.15, 12.7, 0.3, fill_color=COL_LIGHT_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.4, 7.18, 12.5, 0.25,
             [("一句话总结：", 10, COL_TITLE_DARK, True),
              ("Palantir 不是传统软件厂商，而是「决策与运营基础设施」——通过统一业务模型和动作体系，把数据、智能和行动连接起来，", 10, COL_TEXT_DARK, False),
              ("帮助客户在真实世界做准确决策并执行。", 10, COL_TITLE_DARK, True)])

# ====================================================================
# SLIDE 3: TABLE OF CONTENTS
# ====================================================================
def slide_03_toc():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 1, "17", "目录 / Contents")

    toc_items = [
        ("01", "理解 Palantir：放弃简单标签"),
        ("02", "发展历程：能力逐层叠加"),
        ("03", "顶层产品体系：四大平台"),
        ("04", "Foundry：数据运营平台"),
        ("05", "Ontology：核心业务语言"),
        ("06", "Action：从分析走向运营"),
        ("07", "AIP：企业 AI 运行环境"),
        ("08", "Apollo & Rubix：交付底座"),
        ("09", "Global Branching：软件工程延伸"),
        ("10", "Gotham：任务运营系统"),
        ("11", "FDE：产品研发体系"),
        ("12", "Bootcamp：销售与验证"),
        ("13", "商业模式：非纯 SaaS"),
        ("14", "经营数据"),
        ("15", "复合飞轮壁垒"),
        ("16", "约束与挑战"),
        ("17", "最终研究结论"),
    ]
    col1 = toc_items[:9]
    col2 = toc_items[9:]
    for i, (num, title) in enumerate(col1):
        y = 1.7 + i * 0.55
        add_circle(s, 0.5, y + 0.03, 0.32, fill_color=COL_LIGHT_BLUE)
        add_text(s, 0.5, y + 0.07, 0.32, 0.25, num, size=11, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, 1.0, y + 0.05, 5.0, 0.35, title, size=12, color=COL_TEXT_DARK)
    for i, (num, title) in enumerate(col2):
        y = 1.7 + i * 0.55
        add_circle(s, 7.0, y + 0.03, 0.32, fill_color=COL_LIGHT_BLUE)
        add_text(s, 7.0, y + 0.07, 0.32, 0.25, num, size=11, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, 7.5, y + 0.05, 5.0, 0.35, title, size=12, color=COL_TEXT_DARK)

# ====================================================================
# SLIDE 4: 01 理解 Palantir
# ====================================================================
def slide_04_understand():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 2, "17", "01  理解 Palantir：放弃简单标签")

    # Common labels
    add_rect(s, 0.3, 1.6, 6.0, 1.5, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.75)
    add_text(s, 0.5, 1.7, 5.6, 0.3, "常见但不够准确的标签：", size=11, color=COL_RED, bold=True)
    labels = ["大数据分析", "商业智能", "数据中台", "AI软件公司", "政府承包商", "咨询公司", "Ontology 平台"]
    for i, label in enumerate(labels):
        col = i % 4
        row = i // 4
        x = 0.5 + col * 1.35
        y = 2.1 + row * 0.4
        add_rect(s, x, y, 1.25, 0.32, fill_color=COL_BG_WHITE, line_color=COL_RED, line_width=0.5)
        add_text(s, x, y + 0.03, 1.25, 0.25, label, size=9, color=COL_RED, align=PP_ALIGN.CENTER)

    # Real definition
    add_rect(s, 6.6, 1.6, 6.4, 1.5, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.75)
    add_text(s, 6.8, 1.7, 6.0, 0.3, "更准确的定义：", size=11, color=COL_GREEN, bold=True)
    add_text(s, 6.8, 2.05, 6.0, 0.95,
             "帮助组织规模化整合数据、决策与运营。不是只负责把数据汇总起来，"
             "也不是只负责让用户分析数据，而是试图将数据、业务逻辑、软件应用、AI和实际业务行动连接在一个系统内。",
             size=10, color=COL_TEXT_DARK)

    # 7 dimensions
    add_text(s, 0.3, 3.35, 12.7, 0.35, "研究 Palantir 必须同时看 7 个维度：", size=14, color=COL_SECTION_HEADER, bold=True)
    dims = [
        ("①", "如何连接和治理数据", COL_SECTION_HEADER),
        ("②", "如何表达企业业务", COL_ICON_BLUE),
        ("③", "如何把规则引入决策", COL_GREEN),
        ("④", "如何把决策转化为行动", COL_ORANGE),
        ("⑤", "如何在复杂环境中部署", COL_PURPLE),
        ("⑥", "如何进入客户并扩张", COL_TEAL),
        ("⑦", "如何持续运营和维护", COL_RED),
    ]
    box_w = 1.78
    gap = 0.05
    start_x = 0.3
    for i, (num, title, c) in enumerate(dims):
        x = start_x + i * (box_w + gap)
        y = 3.9
        add_rect(s, x, y, box_w, 1.5, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
        add_circle(s, x + box_w/2 - 0.22, y + 0.15, 0.44, fill_color=c)
        add_text(s, x, y + 0.2, box_w, 0.35, num, size=18, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.7, box_w - 0.1, 0.7, title, size=10, color=COL_TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Bottom flow
    add_rect(s, 0.3, 5.7, 12.7, 1.2, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 5.8, 12.3, 0.3, "核心闭环：", size=11, color=COL_SECTION_HEADER, bold=True)
    flow_items = ["感知数据", "发现问题", "生成方案", "权限校验", "执行 Action", "结果反馈", "持续优化"]
    box_w2 = 1.6
    gap2 = 0.17
    start_x2 = 0.5 + (12.3 - box_w2 * 7 - gap2 * 6) / 2
    for i, item in enumerate(flow_items):
        x = start_x2 + i * (box_w2 + gap2)
        y = 6.2
        c = [COL_SECTION_HEADER, COL_ICON_BLUE, COL_GREEN, COL_ORANGE, COL_PURPLE, COL_RED, COL_TEAL][i]
        add_rect(s, x, y, box_w2, 0.5, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_text(s, x, y + 0.1, box_w2, 0.3, item, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        if i < len(flow_items) - 1:
            add_arrow(s, x + box_w2 - 0.02, y + 0.13, gap2 + 0.04, 0.25)

# ====================================================================
# SLIDE 5: 02 发展历程 (was image)
# ====================================================================
def slide_05_history():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 3, "17", "02  发展历程：能力逐层叠加")

    # Timeline
    phases = [
        ("2003-2008", "Gotham 起源", "反恐情报分析\nGovt 市场", COL_RED),
        ("2008-2016", "Palantir Foundry 雏形", "商业市场拓展\n数据整合平台", COL_ORANGE),
        ("2016-2020", "Foundry 成熟", "Ontology 统一模型\nWorkshop + AIP", COL_GREEN),
        ("2020-2023", "AIP 发布", "大模型接入企业\nAI 运行环境", COL_TEAL),
        ("2023-2026", "全面运营基础设施", "Apollo 持续交付\n跨环境部署", COL_PURPLE),
    ]
    # Timeline line
    add_line(s, 0.8, 2.5, 12.5, 2.5, color=COL_BORDER_BLUE, width=2.0)
    box_w = 2.3
    gap = 0.15
    start_x = 0.8
    for i, (period, title, desc, c) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        y = 1.8
        # Circle on timeline
        add_circle(s, x + box_w/2 - 0.15, 2.35, 0.3, fill_color=c)
        add_text(s, x + box_w/2 - 0.15, 2.38, 0.3, 0.25, str(i+1), size=11, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Period label above
        add_text(s, x, y, box_w, 0.3, period, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        # Title
        add_text(s, x, y + 0.3, box_w, 0.35, title, size=11, color=COL_TITLE_DARK, bold=True, align=PP_ALIGN.CENTER)
        # Card below
        add_rect(s, x, 2.9, box_w, 1.6, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_text(s, x + 0.1, 3.0, box_w - 0.2, 1.4, desc, size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Key insight
    add_rect(s, 0.3, 5.0, 12.7, 1.0, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 5.1, 12.3, 0.3, "关键洞察：", size=12, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 0.5, 5.45, 12.3, 0.5,
             "Palantir 不是一步到位的平台，而是从「情报分析」到「数据整合」到「业务建模」到「AI运营」逐步叠加能力层。"
             "每一层不是替代前一层，而是在其基础上增加新的运营维度。",
             size=10, color=COL_TEXT_DARK)

    # Milestones
    add_text(s, 0.3, 6.2, 12.7, 0.3, "关键里程碑：", size=11, color=COL_SECTION_HEADER, bold=True)
    milestones = [
        ("2003", "成立"),
        ("2008", "Gothan 发布"),
        ("2010", "Foundry 雏形"),
        ("2018", "Foundry 商用"),
        ("2020", "直接上市"),
        ("2023", "AIP 发布"),
        ("2025", "$44.75亿收入"),
        ("2026", "954+ 客户"),
    ]
    for i, (year, event) in enumerate(milestones):
        x = 0.3 + i * 1.6
        y = 6.55
        add_rect(s, x, y, 1.45, 0.65, fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_text(s, x, y + 0.05, 1.45, 0.3, year, size=12, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + 0.35, 1.45, 0.25, event, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 6: 03 顶层产品体系 (was image)
# ====================================================================
def slide_06_platforms():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 4, "17", "03  顶层产品体系：四大平台")

    # Central concept: Foundry is the core
    add_circle(s, 5.5, 2.8, 2.0, fill_color=COL_SECTION_HEADER)
    add_text(s, 5.5, 3.15, 2.0, 0.4, "Foundry", size=16, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 5.5, 3.55, 2.0, 0.3, "数据运营平台", size=9, color=COL_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, 5.5, 3.85, 2.0, 0.3, "核心底座", size=8, color=COL_BORDER_BLUE, align=PP_ALIGN.CENTER)

    # Surrounding platforms
    platforms = [
        ("Ontology", "业务运行模型", "对象 · 关系 · Action", COL_GREEN, 1.5, 1.5),
        ("AIP", "AI 运行环境", "模型 · Agent · 工具", COL_PURPLE, 8.5, 1.5),
        ("Apollo", "持续交付底座", "部署 · 升级 · 多环境", COL_ORANGE, 1.5, 5.0),
        ("Rubix", "运行底座", "K8s · 安全 · 扩缩容", COL_TEAL, 8.5, 5.0),
    ]
    for name, role, desc, c, x, y in platforms:
        # Card
        add_rect(s, x, y, 3.3, 1.6, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_plain_rect(s, x, y, 0.08, 1.6, fill_color=c)
        add_text(s, x + 0.25, y + 0.1, 2.0, 0.4, name, size=16, color=c, bold=True)
        add_text(s, x + 0.25, y + 0.55, 2.8, 0.3, role, size=10, color=COL_TEXT_DARK)
        add_text(s, x + 0.25, y + 0.9, 2.8, 0.35, desc, size=9, color=COL_TEXT_GRAY)
        # Arrow to center
        cx, cy = 6.5, 3.8

    # Gotham at bottom
    add_rect(s, 4.5, 5.2, 4.3, 1.3, fill_color=COL_BG_WHITE, line_color=COL_RED, line_width=1.5)
    add_plain_rect(s, 4.5, 5.2, 0.08, 1.3, fill_color=COL_RED)
    add_text(s, 4.75, 5.3, 3.0, 0.4, "Gotham", size=14, color=COL_RED, bold=True)
    add_text(s, 4.75, 5.7, 3.5, 0.3, "任务运营系统（国防/情报）", size=10, color=COL_TEXT_DARK)
    add_text(s, 4.75, 6.05, 3.5, 0.3, "传感器 · 地理空间 · 目标 · 资产调度", size=9, color=COL_TEXT_GRAY)

    # Connecting arrows
    add_arrow(s, 3.8, 2.1, 1.8, 0.2)
    add_arrow(s, 7.8, 2.1, -1.8, 0.2)
    add_arrow(s, 3.8, 5.6, 1.8, 0.2)
    add_arrow(s, 7.8, 5.6, -1.8, 0.2)

    # Bottom: supplementary
    add_text(s, 0.3, 6.8, 12.7, 0.3,
             "Foundry 是核心数据运营平台，Ontology 提供业务语义层，AIP 在其上叠加 AI 能力，Apollo/Rubix 负责持续交付和运行，Gotham 面向国防情报场景。",
             size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 7: 04 Foundry
# ====================================================================
def slide_07_foundry():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 5, "17", "04  Foundry：数据运营平台",
                  "不是数据仓库，而是数据运营平台")

    # Left: Responsibilities
    add_text(s, 0.3, 1.6, 4.0, 0.3, "职责", size=13, color=COL_SECTION_HEADER, bold=True)
    responsibilities = [
        ("数据连接", "多源数据接入"),
        ("数据转换", "ETL / 管道"),
        ("数据质量", "血缘 + 校验"),
        ("数据权限", "字段级控制"),
        ("Ontology 建模", "对象 + 关系"),
        ("业务逻辑", "规则 + 函数"),
        ("模型运行", "ML / LLM"),
        ("应用开发", "Workshop"),
        ("自动化", "工作流 + Action"),
    ]
    for i, (title, sub) in enumerate(responsibilities):
        col = i // 5
        row = i % 5
        x = 0.3 + col * 2.0
        y = 2.0 + row * 0.65
        add_rect(s, x, y, 1.85, 0.55, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
        add_text(s, x + 0.1, y + 0.03, 1.7, 0.25, title, size=10, color=COL_SECTION_HEADER, bold=True)
        add_text(s, x + 0.1, y + 0.28, 1.7, 0.22, sub, size=8, color=COL_TEXT_GRAY)

    # Right: Multimodal data
    add_rect(s, 4.4, 1.6, 4.3, 2.8, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
    add_text(s, 4.6, 1.7, 3.9, 0.3, "Multimodal Data Plane", size=12, color=COL_SECTION_HEADER, bold=True)
    data_types = [
        ("结构化数据", "表格 · 关系数据库", COL_ICON_BLUE),
        ("非结构化", "文档 · 图片 · 视频", COL_GREEN),
        ("流数据", "实时事件 · IoT", COL_ORANGE),
        ("地理空间", "坐标 · 地图 · 卫星", COL_PURPLE),
    ]
    for i, (title, desc, c) in enumerate(data_types):
        y = 2.1 + i * 0.55
        add_circle(s, 4.6, y + 0.03, 0.18, fill_color=c)
        add_text(s, 4.9, y, 3.6, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 4.9, y + 0.22, 3.6, 0.22, desc, size=8, color=COL_TEXT_GRAY)

    # Connect ERP/CRM
    add_text(s, 4.6, 4.35, 3.9, 0.3, "连接：", size=9, color=COL_TEXT_GRAY)
    add_text(s, 4.6, 4.55, 3.9, 0.5,
             "ERP / CRM / 数据仓库 / 工业数据库 / 传感器 / Snowflake / Databricks",
             size=8, color=COL_TEXT_GRAY)

    # Rightmost: "运营统一 ≠ 数据统一"
    add_rect(s, 8.9, 1.6, 4.1, 5.3, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.75)
    add_text(s, 9.1, 1.7, 3.7, 0.35, "运营统一 ≠ 数据统一", size=13, color=COL_SECTION_HEADER, bold=True)

    # Traditional
    add_rect(s, 9.1, 2.2, 3.7, 1.5, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.5)
    add_text(s, 9.2, 2.3, 3.5, 0.25, "传统终点：", size=10, color=COL_RED, bold=True)
    trad_items = ["数据集 / 指标 / 报表", "模型输出"]
    for i, item in enumerate(trad_items):
        add_text(s, 9.3, 2.6 + i * 0.35, 3.3, 0.3, f"• {item}", size=9, color=COL_TEXT_DARK)

    # Foundry
    add_rect(s, 9.1, 3.9, 3.7, 2.8, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.5)
    add_text(s, 9.2, 4.0, 3.5, 0.25, "Foundry 终点：", size=10, color=COL_GREEN, bold=True)
    foundry_items = ["业务对象 / 运营应用", "任务 / 工作流 / Action", "业务系统写回"]
    for i, item in enumerate(foundry_items):
        add_text(s, 9.3, 4.3 + i * 0.35, 3.3, 0.3, f"• {item}", size=9, color=COL_TEXT_DARK)

    # Arrow
    add_down_arrow(s, 10.6, 3.7, 0.6, 0.2, color=COL_ARROW_GRAY)

# ====================================================================
# SLIDE 8: 05 Ontology
# ====================================================================
def slide_08_ontology():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 6, "17", "05  Ontology：核心业务语言",
                  '"Ontology 被设计用于表达企业复杂而相互关联的决策，而不仅仅是表达数据。"')

    # 4 quadrants
    quads = [
        ("Data", "决策依赖的事实", "当前状态 · 历史 · 实时 · 外部", COL_ICON_BLUE, 0.3, 2.0),
        ("Logic", "计算和推理", "规则 · 指标 · ML模型 · LLM函数", COL_GREEN, 6.8, 2.0),
        ("Action", "影响现实世界", "修改对象 · 审批 · 写回系统", COL_ORANGE, 0.3, 4.3),
        ("Security", "权限控制", "查看 · 运行 · 发起 · 修改", COL_PURPLE, 6.8, 4.3),
    ]
    for name, role, desc, c, x, y in quads:
        w = 6.2
        h = 2.1
        add_rect(s, x, y, w, h, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_plain_rect(s, x, y, w, 0.08, fill_color=c)
        # Large label
        add_circle(s, x + 0.2, y + 0.25, 0.6, fill_color=c)
        add_text(s, x + 0.2, y + 0.38, 0.6, 0.35, name[0], size=20, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 1.0, y + 0.2, 2.5, 0.4, name, size=20, color=c, bold=True)
        add_text(s, x + 1.0, y + 0.65, 3.0, 0.3, role, size=11, color=COL_TEXT_DARK)
        add_text(s, x + 1.0, y + 1.0, 5.0, 0.35, desc, size=9, color=COL_TEXT_GRAY)

    # Bottom: noun + verb
    add_rect(s, 0.3, 6.6, 12.7, 0.65, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 6.65, 12.3, 0.25,
             [("名词（对象 + 关系）", 11, COL_SECTION_HEADER, True),
              (" → 描述企业     ", 11, COL_TEXT_DARK, False),
              ("动词（Action）", 11, COL_ORANGE, True),
              (" → 改变企业", 11, COL_TEXT_DARK, False)])
    add_text(s, 0.5, 6.9, 12.3, 0.3,
             "Ontology 不是静态知识图谱，而是企业业务对象运行时，架构的核心系统，主要贯穿 Foundry 和 AIP。",
             size=9, color=COL_TEXT_GRAY)

# ====================================================================
# SLIDE 9: Ontology detail (was image)
# ====================================================================
def slide_09_ontology_detail():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 6, "17", "05  Ontology：四大要素详解")

    elements = [
        ("Object Types\n对象类型", "企业实体的数字化表达", [
            "Customer · Order · Product · Factory",
            "属性（Properties）：名称、金额、状态",
            "主键 + 主时间戳",
            "支持继承和接口",
        ], COL_ICON_BLUE),
        ("Link Types\n关系类型", "对象之间的连接", [
            "Customer → placed → Order",
            "Order → contains → Product",
            "Cardinality: one-to-many / many-to-many",
            "方向性 · 可带属性",
        ], COL_GREEN),
        ("Action Types\n动作类型", "改变现实世界的操作", [
            "输入参数 + 表单",
            "业务规则 + 提交条件",
            "权限控制",
            "对象修改 + 外部写回",
        ], COL_ORANGE),
        ("Functions\n函数类型", "业务逻辑的封装", [
            "计算函数（利润率）",
            "决策函数（信用评估）",
            "ML/LLM 模型调用",
            "可被 Action 和 Workshop 调用",
        ], COL_PURPLE),
    ]
    for i, (title, subtitle, items, c) in enumerate(elements):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.4
        y = 1.6 + row * 2.7
        w = 6.2
        h = 2.5
        add_rect(s, x, y, w, h, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 0.08, h, fill_color=c)
        add_text(s, x + 0.25, y + 0.1, 3.0, 0.5, title, size=14, color=c, bold=True)
        add_text(s, x + 0.25, y + 0.65, 5.5, 0.25, subtitle, size=10, color=COL_TEXT_DARK)
        for j, item in enumerate(items):
            add_circle(s, x + 0.3, y + 1.0 + j * 0.32 + 0.03, 0.06, fill_color=c)
            add_text(s, x + 0.45, y + 1.0 + j * 0.32, 5.3, 0.3, item, size=9, color=COL_TEXT_GRAY)

# ====================================================================
# SLIDE 10: Action feedback loop (was image)
# ====================================================================
def slide_10_action_loop():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 7, "17", "05  Ontology：决策反馈循环")

    # Central circle
    cx, cy, cd = 5.0, 3.5, 2.5
    add_circle(s, cx, cy, cd, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=2.0)
    add_text(s, cx, cy + 0.3, cd, 0.5, "Ontology", size=22, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy + 0.9, cd, 0.3, "统一业务模型", size=12, color=COL_SUB_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy + 1.3, cd, 0.3, "对象·关系·Action", size=10, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Surrounding nodes
    nodes = [
        ("Data\n数据", COL_ICON_BLUE, 3.0, 1.0),
        ("Logic\n逻辑", COL_GREEN, 7.5, 1.0),
        ("Action\n行动", COL_ORANGE, 3.0, 5.5),
        ("Security\n安全", COL_PURPLE, 7.5, 5.5),
    ]
    for name, c, nx, ny in nodes:
        add_rect(s, nx, ny, 2.0, 1.0, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_text(s, nx, ny + 0.15, 2.0, 0.7, name, size=13, color=c, bold=True, align=PP_ALIGN.CENTER)

    # Right side: feedback loop description
    add_rect(s, 9.8, 1.6, 3.2, 5.0, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 10.0, 1.7, 2.8, 0.3, "决策反馈循环：", size=12, color=COL_SECTION_HEADER, bold=True)
    feedback_items = [
        "Action 结果", "→ 重新进入 Ontology", "→ 调整规则", "→ 重训模型", "→ 优化建议", "→ 更好的决策",
    ]
    for i, item in enumerate(feedback_items):
        y = 2.1 + i * 0.6
        add_text(s, 10.1, y, 2.7, 0.35, item, size=10, color=COL_TEXT_DARK if i == 0 else COL_TEXT_GRAY,
                 bold=(i == 0))
        if i < len(feedback_items) - 1:
            add_down_arrow(s, 11.2, y + 0.35, 0.15, 0.15, color=COL_ARROW_GRAY)

    # Bottom note
    add_text(s, 0.3, 6.8, 9.0, 0.3,
             "Data → Logic → Action → 结果反馈 → 回到 Ontology → 持续优化",
             size=10, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 11: Action types (was image)
# ====================================================================
def slide_11_action_types():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 7, "17", "05  Ontology：Action 类型与反馈")

    # Left: Action components
    add_rect(s, 0.3, 1.6, 6.0, 5.3, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
    add_text(s, 0.5, 1.7, 5.6, 0.3, "Action 包含：", size=13, color=COL_SECTION_HEADER, bold=True)
    action_parts = [
        ("输入参数 + 表单", "用户交互界面", COL_ICON_BLUE),
        ("业务规则 + 提交条件", "验证逻辑", COL_GREEN),
        ("权限控制", "谁能执行", COL_PURPLE),
        ("对象修改", "一个事务中修改多个对象", COL_ORANGE),
        ("外部副作用", "写回业务系统（ERP/CRM）", COL_RED),
        ("审计 + 运行指标", "可追溯、可度量", COL_TEAL),
    ]
    for i, (title, desc, c) in enumerate(action_parts):
        y = 2.15 + i * 0.75
        add_circle(s, 0.6, y + 0.08, 0.28, fill_color=c)
        add_text(s, 0.6, y + 0.12, 0.28, 0.25, str(i+1), size=11, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, 1.05, y + 0.02, 4.8, 0.3, title, size=11, color=COL_TITLE_DARK, bold=True)
        add_text(s, 1.05, y + 0.32, 4.8, 0.25, desc, size=9, color=COL_TEXT_GRAY)

    # Right: Traditional vs Palantir
    add_rect(s, 6.6, 1.6, 6.4, 2.3, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.75)
    add_text(s, 6.8, 1.7, 6.0, 0.3, "传统数据/AI 平台流程：", size=12, color=COL_RED, bold=True)
    trad_flow = ["发现问题", "提供建议", "用户离开平台执行"]
    for i, item in enumerate(trad_flow):
        x = 6.9 + i * 2.0
        add_rect(s, x, 2.15, 1.7, 0.5, fill_color=COL_BG_WHITE, line_color=COL_RED, line_width=0.75)
        add_text(s, x, 2.25, 1.7, 0.3, item, size=10, color=COL_RED, align=PP_ALIGN.CENTER)
        if i < len(trad_flow) - 1:
            add_text(s, x + 1.65, 2.25, 0.4, 0.3, "→", size=14, color=COL_RED, align=PP_ALIGN.CENTER)
    add_text(s, 6.8, 2.9, 6.0, 0.3, "→ 决策与执行断裂", size=10, color=COL_RED, bold=True)
    add_text(s, 6.8, 3.2, 6.0, 0.3, "→ 无法追踪执行结果", size=10, color=COL_RED)

    add_rect(s, 6.6, 4.1, 6.4, 2.8, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.75)
    add_text(s, 6.8, 4.2, 6.0, 0.3, "Palantir 流程：", size=12, color=COL_GREEN, bold=True)
    pal_flow = ["发现问题", "生成方案", "权限校验", "执行 Action", "结果反馈"]
    bw = 1.15
    gap = 0.08
    for i, item in enumerate(pal_flow):
        x = 6.9 + i * (bw + gap)
        add_rect(s, x, 4.65, bw, 0.5, fill_color=COL_BG_WHITE, line_color=COL_GREEN, line_width=0.75)
        add_text(s, x, 4.75, bw, 0.3, item, size=8, color=COL_GREEN, align=PP_ALIGN.CENTER)
        if i < len(pal_flow) - 1:
            add_text(s, x + bw - 0.05, 4.75, 0.2, 0.3, "→", size=10, color=COL_GREEN, align=PP_ALIGN.CENTER)
    add_text(s, 6.8, 5.4, 6.0, 0.3, "→ 闭环：决策 → 执行 → 反馈", size=10, color=COL_GREEN, bold=True)
    add_text(s, 6.8, 5.7, 6.0, 0.3, "→ Action 结果重新进入 Ontology", size=10, color=COL_GREEN)
    add_text(s, 6.8, 6.0, 6.0, 0.3, "→ 调整规则 / 重训模型 / 优化建议", size=10, color=COL_GREEN)
    add_text(s, 6.8, 6.3, 6.0, 0.3, "→ 形成持续改进的运营闭环", size=10, color=COL_GREEN, bold=True)

# ====================================================================
# SLIDE 12: 06 Action
# ====================================================================
def slide_12_action():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 8, "17", "06  Action：从分析走向运营的关键")

    # Top comparison
    # Traditional
    add_rect(s, 0.3, 1.6, 6.0, 1.2, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.75)
    add_text(s, 0.5, 1.65, 5.6, 0.25, "传统数据/AI 平台流程", size=10, color=COL_RED, bold=True)
    trad_steps = ["发现问题", "提供建议", "用户离开平台执行"]
    for i, step in enumerate(trad_steps):
        x = 0.5 + i * 1.9
        add_rect(s, x, 2.0, 1.7, 0.55, fill_color=COL_BG_WHITE, line_color=COL_RED, line_width=0.5)
        add_text(s, x, 2.12, 1.7, 0.3, step, size=9, color=COL_RED, align=PP_ALIGN.CENTER)
        if i < len(trad_steps) - 1:
            add_arrow(s, x + 1.68, 2.18, 0.22, 0.2, color=COL_RED)

    # Palantir
    add_rect(s, 6.6, 1.6, 6.4, 1.2, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.75)
    add_text(s, 6.8, 1.65, 6.0, 0.25, "Palantir 流程", size=10, color=COL_GREEN, bold=True)
    pal_steps = ["发现问题", "生成方案", "权限校验", "执行 Action", "结果反馈"]
    bw = 1.15
    gap = 0.1
    for i, step in enumerate(pal_steps):
        x = 6.8 + i * (bw + gap)
        add_rect(s, x, 2.0, bw, 0.55, fill_color=COL_BG_WHITE, line_color=COL_GREEN, line_width=0.5)
        add_text(s, x, 2.12, bw, 0.3, step, size=8, color=COL_GREEN, align=PP_ALIGN.CENTER)
        if i < len(pal_steps) - 1:
            add_arrow(s, x + bw - 0.02, 2.18, gap + 0.04, 0.2, color=COL_GREEN)

    # Middle: Action contains
    add_text(s, 0.3, 3.0, 12.7, 0.3, "Action 包含：", size=13, color=COL_SECTION_HEADER, bold=True)
    action_items = [
        ("输入参数", "表单界面", COL_ICON_BLUE),
        ("业务规则", "提交条件验证", COL_GREEN),
        ("权限控制", "谁能发起", COL_PURPLE),
        ("对象修改", "事务性多对象修改", COL_ORANGE),
        ("外部副作用", "写回业务系统", COL_RED),
        ("审计指标", "可追溯、可度量", COL_TEAL),
    ]
    for i, (title, desc, c) in enumerate(action_items):
        x = 0.3 + i * 2.13
        y = 3.4
        add_rect(s, x, y, 2.0, 1.3, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 2.0, 0.06, fill_color=c)
        add_text(s, x + 0.05, y + 0.15, 1.9, 0.3, title, size=11, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.5, 1.9, 0.3, desc, size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Bottom: Feedback loop
    add_rect(s, 0.3, 5.0, 12.7, 1.9, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 5.1, 12.3, 0.3, "决策反馈循环：", size=12, color=COL_SECTION_HEADER, bold=True)
    loop_items = [
        ("Action 执行", COL_ORANGE),
        ("→ 结果进入 Ontology", COL_ICON_BLUE),
        ("→ 调整业务规则", COL_GREEN),
        ("→ 重训 ML 模型", COL_PURPLE),
        ("→ 优化决策建议", COL_RED),
        ("→ 更好的 Action", COL_TEAL),
    ]
    bw2 = 1.9
    gap2 = 0.13
    start_x2 = 0.5 + (12.3 - bw2 * 6 - gap2 * 5) / 2
    for i, (item, c) in enumerate(loop_items):
        x = start_x2 + i * (bw2 + gap2)
        y = 5.6
        add_rect(s, x, y, bw2, 0.65, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_text(s, x, y + 0.15, bw2, 0.35, item, size=9, color=c, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, 0.5, 6.45, 12.3, 0.3,
             "语义 + 动作 + 反馈 = 从「分析工具」升级为「运营系统」",
             size=11, color=COL_TITLE_DARK, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 13: 07 AIP
# ====================================================================
def slide_13_aip():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 9, "17", "07  AIP：企业 AI 运行环境",
                  "不是大模型网关，而是企业 AI 运行环境")

    # Left: Without Foundry + Ontology
    add_rect(s, 0.3, 1.6, 4.0, 2.8, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.75)
    add_text(s, 0.5, 1.7, 3.6, 0.3, "没有 Foundry + Ontology", size=11, color=COL_RED, bold=True)
    add_text(s, 0.5, 2.05, 3.6, 0.25, "大模型只能看到：", size=9, color=COL_TEXT_DARK)
    without_items = ["数据库字段", "搜索结果", "文档片段", "Prompt 临时上下文"]
    for i, item in enumerate(without_items):
        add_text(s, 0.7, 2.4 + i * 0.35, 3.3, 0.3, f"• {item}", size=9, color=COL_TEXT_GRAY)
    add_text(s, 0.5, 3.85, 3.6, 0.3, "→ AI 无法理解业务上下文", size=9, color=COL_RED, bold=True)

    # Arrow
    add_arrow(s, 4.3, 2.8, 0.4, 0.3, color=COL_SECTION_HEADER)

    # Middle: Ontology bridge
    add_rect(s, 4.8, 1.6, 3.7, 2.8, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=1.5)
    add_text(s, 5.0, 1.8, 3.3, 0.35, "接入 Ontology 后", size=11, color=COL_SECTION_HEADER, bold=True)
    with_items = [
        ("当前对象是什么 + 关系", COL_ICON_BLUE),
        ("可调用的规则和模型", COL_GREEN),
        ("可执行的 Action", COL_ORANGE),
        ("当前用户拥有的权限", COL_PURPLE),
    ]
    for i, (item, c) in enumerate(with_items):
        add_circle(s, 5.1, 2.3 + i * 0.45 + 0.03, 0.1, fill_color=c)
        add_text(s, 5.3, 2.3 + i * 0.45, 3.0, 0.35, item, size=9, color=COL_TEXT_DARK)

    # Arrow
    add_arrow(s, 8.55, 2.8, 0.4, 0.3, color=COL_SECTION_HEADER)

    # Right: AIP core capabilities
    add_rect(s, 9.1, 1.6, 3.9, 2.8, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.75)
    add_text(s, 9.3, 1.7, 3.5, 0.3, "AIP 核心能力", size=11, color=COL_GREEN, bold=True)
    caps = [
        "多模型连接 + Agent 开发",
        "自动化 + AI 应用",
        "上下文工程 + 工具调用",
        "AI 评测 + 可观测",
        "安全治理 + 人工审核",
    ]
    for i, cap in enumerate(caps):
        add_circle(s, 9.3, 2.1 + i * 0.4 + 0.03, 0.1, fill_color=COL_GREEN)
        add_text(s, 9.5, 2.1 + i * 0.4, 3.3, 0.35, cap, size=9, color=COL_TEXT_DARK)

    # Bottom: Agent inherits permissions
    add_rect(s, 0.3, 4.7, 12.7, 2.2, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 4.8, 12.3, 0.3, "Agent 继承现实权限", size=13, color=COL_SECTION_HEADER, bold=True)

    agent_flow = [
        ("用户权限", "身份 + 角色", COL_ICON_BLUE),
        ("→ Agent", "继承权限", COL_SECTION_HEADER),
        ("→ 调用 Ontology", "只能看到有权限的对象", COL_GREEN),
        ("→ 执行 Action", "在权限范围内操作", COL_ORANGE),
        ("→ 审计日志", "全程可追溯", COL_PURPLE),
    ]
    bw = 2.35
    gap = 0.1
    start_x = 0.5 + (12.3 - bw * 5 - gap * 4) / 2
    for i, (title, desc, c) in enumerate(agent_flow):
        x = start_x + i * (bw + gap)
        y = 5.3
        add_rect(s, x, y, bw, 1.2, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_text(s, x + 0.05, y + 0.1, bw - 0.1, 0.35, title, size=11, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.5, bw - 0.1, 0.6, desc, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    add_text(s, 0.3, 6.7, 12.7, 0.3,
             "AIP 不是在 AI 之上加一层企业外壳，而是在企业运营系统之上叠加 AI 能力 — 让 AI 理解业务语义和权限边界。",
             size=10, color=COL_TEXT_DARK, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 14: AIP Architecture (was image)
# ====================================================================
def slide_14_aip_arch():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 9, "17", "07  AIP：架构全景")

    # Layer 1: Models
    add_text(s, 0.3, 1.6, 2.5, 0.3, "模型层", size=11, color=COL_PURPLE, bold=True)
    models = ["GPT-4o", "Claude", "Gemini", "GLM-4", "Llama", "自定义"]
    for i, m in enumerate(models):
        x = 0.3 + i * 1.3
        add_rect(s, x, 1.95, 1.2, 0.45, fill_color=COL_PURPLE_BG, line_color=COL_PURPLE, line_width=0.5)
        add_text(s, x, 2.05, 1.2, 0.25, m, size=9, color=COL_PURPLE, align=PP_ALIGN.CENTER)

    add_down_arrow(s, 6.5, 2.45, 0.3, 0.2)

    # Layer 2: AIP Core
    add_rect(s, 0.3, 2.8, 12.7, 1.8, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=1.5)
    add_text(s, 0.5, 2.9, 5.0, 0.3, "AIP 核心引擎", size=13, color=COL_SECTION_HEADER, bold=True)
    core_items = [
        ("模型路由", "场景化选模", COL_ICON_BLUE),
        ("Agent 编排", "多步推理", COL_GREEN),
        ("工具调用", "Query/Function/Action", COL_ORANGE),
        ("上下文工程", "Ontology 注入", COL_PURPLE),
        ("安全护栏", "权限 + 出境", COL_RED),
        ("评测可观测", "Eval + Trace", COL_TEAL),
    ]
    for i, (title, desc, c) in enumerate(core_items):
        x = 0.5 + i * 2.05
        y = 3.3
        add_rect(s, x, y, 1.95, 1.1, fill_color=COL_BG_WHITE, line_color=c, line_width=0.75)
        add_text(s, x + 0.05, y + 0.1, 1.85, 0.3, title, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.45, 1.85, 0.3, desc, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    add_down_arrow(s, 6.5, 4.65, 0.3, 0.2)

    # Layer 3: Ontology + Foundry
    add_rect(s, 0.3, 5.0, 12.7, 1.5, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=1.5)
    add_text(s, 0.5, 5.1, 5.0, 0.3, "Ontology + Foundry（业务运行时）", size=13, color=COL_GREEN, bold=True)
    onto_items = [
        "Object Types（对象）",
        "Link Types（关系）",
        "Action Types（动作）",
        "Functions（函数）",
        "Datasets（数据集）",
        "Permissions（权限）",
    ]
    for i, item in enumerate(onto_items):
        x = 0.5 + i * 2.05
        add_rect(s, x, 5.55, 1.95, 0.75, fill_color=COL_BG_WHITE, line_color=COL_GREEN, line_width=0.5)
        add_text(s, x, 5.7, 1.95, 0.4, item, size=9, color=COL_GREEN, align=PP_ALIGN.CENTER)

    # Bottom note
    add_text(s, 0.3, 6.7, 12.7, 0.3,
             "AIP 的模型层 → AIP 核心引擎（路由/编排/工具/安全）→ Ontology + Foundry（业务运行时）— 三层贯穿，AI 始终在业务上下文中运行。",
             size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 15: AIP Agent (was image)
# ====================================================================
def slide_15_aip_agent():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 9, "17", "07  AIP：Agent 运行模型")

    # Center: Agent
    add_circle(s, 5.4, 2.8, 2.2, fill_color=COL_SECTION_HEADER)
    add_text(s, 5.4, 3.1, 2.2, 0.4, "Agent", size=22, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 5.4, 3.6, 2.2, 0.3, "继承用户权限", size=11, color=COL_LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Left: Inputs
    add_text(s, 0.3, 1.6, 3.0, 0.3, "输入", size=12, color=COL_ICON_BLUE, bold=True)
    inputs = [
        ("自然语言", "用户意图"),
        ("Ontology 查询", "对象+关系"),
        ("业务规则", "Functions"),
        ("权限上下文", "角色+Marking"),
    ]
    for i, (title, desc) in enumerate(inputs):
        y = 2.0 + i * 0.75
        add_rect(s, 0.3, y, 2.8, 0.65, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
        add_text(s, 0.4, y + 0.05, 2.6, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 0.4, y + 0.32, 2.6, 0.25, desc, size=8, color=COL_TEXT_GRAY)

    add_arrow(s, 3.1, 3.5, 0.4, 0.3, color=COL_SECTION_HEADER)

    # Right: Outputs / Tools
    add_text(s, 9.5, 1.6, 3.5, 0.3, "工具 / 输出", size=12, color=COL_ORANGE, bold=True)
    tools = [
        ("Query Objects", "查询业务对象"),
        ("Apply Function", "调用计算/规则"),
        ("Execute Action", "执行业务动作"),
        ("Write Back", "写回外部系统"),
    ]
    for i, (title, desc) in enumerate(tools):
        y = 2.0 + i * 0.75
        add_rect(s, 9.5, y, 3.5, 0.65, fill_color=COL_ORANGE_BG, line_color=COL_ORANGE, line_width=0.5)
        add_text(s, 9.6, y + 0.05, 3.3, 0.25, title, size=10, color=COL_ORANGE, bold=True)
        add_text(s, 9.6, y + 0.32, 3.3, 0.25, desc, size=8, color=COL_TEXT_GRAY)

    add_arrow(s, 7.6, 3.5, 0.4, 0.3, color=COL_SECTION_HEADER)

    # Bottom: Guardrails
    add_rect(s, 0.3, 5.2, 12.7, 1.7, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.75)
    add_text(s, 0.5, 5.3, 12.3, 0.3, "安全护栏（Guardrails）", size=12, color=COL_RED, bold=True)
    guards = [
        ("权限继承", "Agent 只能做用户能做的事"),
        ("数据出境", "PII 数据禁止调用境外模型"),
        ("人工审核", "高风险操作需人工确认"),
        ("审计日志", "全链路可追溯"),
    ]
    for i, (title, desc) in enumerate(guards):
        x = 0.5 + i * 3.1
        add_rect(s, x, 5.7, 2.9, 1.0, fill_color=COL_BG_WHITE, line_color=COL_RED, line_width=0.5)
        add_text(s, x + 0.1, 5.78, 2.7, 0.3, title, size=10, color=COL_RED, bold=True)
        add_text(s, x + 0.1, 6.1, 2.7, 0.35, desc, size=8, color=COL_TEXT_GRAY)

# ====================================================================
# SLIDE 16: 08 Apollo & Rubix
# ====================================================================
def slide_16_apollo():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 10, "17", "08  Apollo & Rubix：经常被低估的壁垒")

    # Apollo
    add_rect(s, 0.3, 1.6, 6.0, 4.5, fill_color=COL_BG_WHITE, line_color=COL_ORANGE, line_width=1.5)
    add_plain_rect(s, 0.3, 1.6, 6.0, 0.08, fill_color=COL_ORANGE)
    add_text(s, 0.5, 1.75, 5.6, 0.35, "Apollo — 持续交付", size=16, color=COL_ORANGE, bold=True)
    add_text(s, 0.5, 2.15, 5.6, 0.3, "每天协调数千次跨数百项服务的零停机升级", size=10, color=COL_TEXT_GRAY)
    apollo_items = [
        ("软件版本管理", "多版本并存、灰度控制"),
        ("配置管理", "多环境配置同步"),
        ("安全更新", "零停机安全补丁"),
        ("灰度部署", "按比例、按客户逐步推出"),
        ("回滚", "秒级回滚到任意版本"),
        ("持续运维", "7×24 自动化运维"),
    ]
    for i, (title, desc) in enumerate(apollo_items):
        y = 2.6 + i * 0.58
        add_circle(s, 0.6, y + 0.05, 0.15, fill_color=COL_ORANGE)
        add_text(s, 0.85, y, 5.2, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 0.85, y + 0.24, 5.2, 0.25, desc, size=8, color=COL_TEXT_GRAY)

    # Rubix
    add_rect(s, 6.6, 1.6, 6.4, 4.5, fill_color=COL_BG_WHITE, line_color=COL_TEAL, line_width=1.5)
    add_plain_rect(s, 6.6, 1.6, 6.4, 0.08, fill_color=COL_TEAL)
    add_text(s, 6.8, 1.75, 6.0, 0.35, "Rubix — 运行底座", size=16, color=COL_TEAL, bold=True)
    add_text(s, 6.8, 2.15, 6.0, 0.3, "强化后的 Kubernetes 运行底座", size=10, color=COL_TEXT_GRAY)
    rubix_items = [
        ("高可用", "多 AZ 自动故障转移"),
        ("自动扩缩容", "按负载弹性伸缩"),
        ("工作负载隔离", "命名空间级隔离"),
        ("默认安全网络", "Zero Trust 网络模型"),
        ("加密 + 身份验证", "全链路 mTLS"),
        ("临时计算节点", "按需创建/销毁"),
        ("多环境一致", "开发/测试/生产一致"),
    ]
    for i, (title, desc) in enumerate(rubix_items):
        y = 2.6 + i * 0.5
        add_circle(s, 6.9, y + 0.05, 0.15, fill_color=COL_TEAL)
        add_text(s, 7.15, y, 5.5, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 7.15, y + 0.24, 5.5, 0.25, desc, size=8, color=COL_TEXT_GRAY)

    # Bottom: Why this matters
    add_rect(s, 0.3, 6.3, 12.7, 0.8, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 6.35, 12.3, 0.3, "为什么这一层重要：", size=11, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 0.5, 6.65, 12.3, 0.35,
             "企业客户需要软件长期运行、持续升级、满足安全要求、适配已有基础设施——Apollo 和 Rubix 解决的是「从可以开发」到「可以长期运营」的跨越。",
             size=9, color=COL_TEXT_DARK)

# ====================================================================
# SLIDE 17: Apollo detail (was image)
# ====================================================================
def slide_17_apollo_detail():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 10, "17", "08  Apollo：持续交付详解")

    # Pipeline visualization
    add_text(s, 0.3, 1.6, 12.7, 0.3, "发布管道", size=13, color=COL_SECTION_HEADER, bold=True)

    pipeline = [
        ("代码提交", "Git Push", COL_ICON_BLUE),
        ("自动构建", "CI Pipeline", COL_GREEN),
        ("测试验证", "Unit + Integration", COL_TEAL),
        ("安全扫描", "Vulnerability Check", COL_ORANGE),
        ("灰度部署", "Canary 1% → 10%", COL_PURPLE),
        ("全量发布", "100% Rollout", COL_RED),
        ("监控验证", "Health Check", COL_SECTION_HEADER),
    ]
    bw = 1.7
    gap = 0.1
    start_x = 0.3 + (12.7 - bw * 7 - gap * 6) / 2
    for i, (title, desc, c) in enumerate(pipeline):
        x = start_x + i * (bw + gap)
        y = 2.0
        add_rect(s, x, y, bw, 1.0, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, bw, 0.06, fill_color=c)
        add_text(s, x + 0.05, y + 0.12, bw - 0.1, 0.3, title, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.5, bw - 0.1, 0.4, desc, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            add_arrow(s, x + bw - 0.02, y + 0.35, gap + 0.04, 0.25, color=COL_ARROW_GRAY)

    # Rollback
    add_rect(s, 2.0, 3.3, 9.3, 0.5, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.5)
    add_text(s, 2.2, 3.35, 8.9, 0.35,
             "← 回滚：任何阶段失败，自动回滚到上一个稳定版本（秒级）",
             size=10, color=COL_RED, bold=True, align=PP_ALIGN.CENTER)

    # Features grid
    add_text(s, 0.3, 4.1, 12.7, 0.3, "Apollo 核心特性", size=13, color=COL_SECTION_HEADER, bold=True)
    features = [
        ("多目标部署", "同时部署到多个环境/区域/客户", COL_ORANGE),
        ("版本管理", "管理数千个服务的版本矩阵", COL_ICON_BLUE),
        ("配置热更新", "不重启更新运行时配置", COL_GREEN),
        ("健康监控", "实时监控部署后系统健康", COL_RED),
        ("灰度策略", "按客户/区域/比例逐步推出", COL_PURPLE),
        ("依赖编排", "服务间依赖关系自动排序", COL_TEAL),
    ]
    for i, (title, desc, c) in enumerate(features):
        col = i % 3
        row = i // 3
        x = 0.3 + col * 4.2
        y = 4.5 + row * 1.3
        add_rect(s, x, y, 4.0, 1.15, fill_color=COL_BG_WHITE, line_color=c, line_width=0.75)
        add_plain_rect(s, x, y, 0.08, 1.15, fill_color=c)
        add_text(s, x + 0.25, y + 0.1, 3.5, 0.3, title, size=11, color=c, bold=True)
        add_text(s, x + 0.25, y + 0.45, 3.5, 0.6, desc, size=9, color=COL_TEXT_GRAY)

    # Stats
    add_text(s, 0.3, 7.0, 12.7, 0.25,
             "规模：每天数千次部署 · 数百项服务 · 零停机升级 · 99.99% 可用性",
             size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 18: 09 Global Branching
# ====================================================================
def slide_18_branching():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 11, "17", "09  Global Branching：软件工程延伸到业务系统")

    # Left: Traditional Git
    add_rect(s, 0.3, 1.6, 4.0, 2.5, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.75)
    add_text(s, 0.5, 1.7, 3.6, 0.3, "传统 Git 管理", size=12, color=COL_RED, bold=True)
    add_text(s, 0.5, 2.1, 3.6, 0.3, "仅管理代码", size=10, color=COL_TEXT_DARK)
    trad_git = ["源代码", "配置文件", "基础设施代码"]
    for i, item in enumerate(trad_git):
        add_text(s, 0.7, 2.5 + i * 0.4, 3.3, 0.35, f"• {item}", size=10, color=COL_TEXT_GRAY)
    add_text(s, 0.5, 3.6, 3.6, 0.3, "→ 业务逻辑无法版本化", size=9, color=COL_RED, bold=True)

    # Arrow
    add_arrow(s, 4.3, 2.5, 0.4, 0.3, color=COL_SECTION_HEADER)

    # Middle: Global Branching
    add_rect(s, 4.8, 1.6, 3.7, 2.5, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.75)
    add_text(s, 5.0, 1.7, 3.3, 0.3, "Global Branching", size=12, color=COL_GREEN, bold=True)
    add_text(s, 5.0, 2.1, 3.3, 0.3, "管理全部", size=10, color=COL_TEXT_DARK)
    gb_items = [
        "数据管道", "数据 Schema", "Ontology", "Action",
        "应用 + Functions", "自动化流程"
    ]
    for i, item in enumerate(gb_items):
        col = i % 2
        row = i // 2
        x = 5.0 + col * 1.6
        y = 2.5 + row * 0.4
        add_rect(s, x, y, 1.5, 0.3, fill_color=COL_BG_WHITE, line_color=COL_GREEN, line_width=0.5)
        add_text(s, x, y + 0.03, 1.5, 0.25, item, size=8, color=COL_GREEN, align=PP_ALIGN.CENTER)

    # Arrow
    add_arrow(s, 8.55, 2.5, 0.4, 0.3, color=COL_SECTION_HEADER)

    # Right: Flow
    add_rect(s, 9.1, 1.6, 3.9, 2.5, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=0.75)
    add_text(s, 9.3, 1.7, 3.5, 0.3, "隔离分支开发流程", size=11, color=COL_SECTION_HEADER, bold=True)
    flow = ["创建分支", "修改完整工作流", "测试（Action可运行但不写主分支）", "审核通过", "统一合并"]
    for i, item in enumerate(flow):
        add_circle(s, 9.3, 2.15 + i * 0.38 + 0.03, 0.1, fill_color=COL_SECTION_HEADER)
        add_text(s, 9.55, 2.15 + i * 0.38, 3.2, 0.35, item, size=8, color=COL_TEXT_DARK)

    # Bottom: Why
    add_rect(s, 0.3, 4.4, 12.7, 2.6, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 4.5, 12.3, 0.35, "为什么需要 Global Branching？", size=13, color=COL_SECTION_HEADER, bold=True)

    reasons = [
        ("业务定义变化", "Ontology 对象/属性/关系的修改需要版本管理", COL_ICON_BLUE),
        ("决策流程变化", "Action 和 Functions 的修改需要测试和审核", COL_GREEN),
        ("安全合规", "敏感修改需要审计和审批流程", COL_PURPLE),
        ("团队协作", "多个团队同时修改不同部分而不互相干扰", COL_ORANGE),
    ]
    for i, (title, desc, c) in enumerate(reasons):
        x = 0.5 + i * 3.1
        y = 5.0
        add_rect(s, x, y, 2.9, 1.8, fill_color=COL_BG_WHITE, line_color=c, line_width=0.75)
        add_plain_rect(s, x, y, 2.9, 0.06, fill_color=c)
        add_text(s, x + 0.1, y + 0.15, 2.7, 0.35, title, size=11, color=c, bold=True)
        add_text(s, x + 0.1, y + 0.55, 2.7, 0.8, desc, size=9, color=COL_TEXT_GRAY)

    add_text(s, 0.3, 7.0, 12.7, 0.3,
             "企业不仅需要管理代码变化，还需要管理业务定义和决策流程的变化。",
             size=10, color=COL_TEXT_DARK, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 19: Global Branching detail (was image)
# ====================================================================
def slide_19_branching_detail():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 11, "17", "09  Global Branching：分支管理全景")

    # Branch tree visualization
    add_text(s, 0.3, 1.6, 12.7, 0.3, "分支模型", size=13, color=COL_SECTION_HEADER, bold=True)

    # Main branch
    add_rect(s, 5.0, 2.1, 3.3, 0.5, fill_color=COL_SECTION_HEADER)
    add_text(s, 5.0, 2.18, 3.3, 0.35, "main（主分支）", size=11, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Feature branches
    branches = [
        ("feature/order-v2", "新增订单 Action", COL_GREEN, 1.0, 3.0),
        ("feature/credit-rule", "修改信用规则", COL_ORANGE, 5.0, 3.0),
        ("feature/new-ot", "新增 Supplier OT", COL_PURPLE, 9.0, 3.0),
    ]
    for name, desc, c, x, y in branches:
        # Branch line
        add_line(s, 6.65, 2.6, x + 1.65, y, color=c, width=1.5)
        add_rect(s, x, y, 3.3, 0.5, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_text(s, x, y + 0.08, 3.3, 0.35, name, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + 0.55, 3.3, 0.25, desc, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # What's in a branch
    add_text(s, 0.3, 4.0, 12.7, 0.3, "一个分支包含什么：", size=12, color=COL_SECTION_HEADER, bold=True)
    branch_items = [
        ("数据管道", "Pipeline Builder 变更", COL_ICON_BLUE),
        ("Schema", "Dataset Schema 变更", COL_GREEN),
        ("Ontology", "OT/Property/Link/Action 变更", COL_ORANGE),
        ("应用", "Workshop 应用变更", COL_PURPLE),
        ("Functions", "函数逻辑变更", COL_RED),
        ("自动化", "Automate 规则变更", COL_TEAL),
    ]
    for i, (title, desc, c) in enumerate(branch_items):
        x = 0.3 + i * 2.13
        y = 4.4
        add_rect(s, x, y, 2.0, 1.2, fill_color=COL_BG_WHITE, line_color=c, line_width=0.75)
        add_text(s, x + 0.05, y + 0.1, 1.9, 0.3, title, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.45, 1.9, 0.6, desc, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Merge process
    add_text(s, 0.3, 5.9, 12.7, 0.3, "合并流程", size=12, color=COL_SECTION_HEADER, bold=True)
    merge_steps = ["创建分支", "隔离修改", "运行测试", "代码审查", "安全审核", "合并到 main", "Apollo 部署"]
    bw = 1.7
    gap = 0.1
    start_x = 0.3 + (12.7 - bw * 7 - gap * 6) / 2
    for i, step in enumerate(merge_steps):
        x = start_x + i * (bw + gap)
        y = 6.3
        c = [COL_ICON_BLUE, COL_GREEN, COL_TEAL, COL_ORANGE, COL_RED, COL_SECTION_HEADER, COL_PURPLE][i]
        add_rect(s, x, y, bw, 0.55, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_text(s, x, y + 0.12, bw, 0.3, step, size=9, color=c, bold=True, align=PP_ALIGN.CENTER)
        if i < len(merge_steps) - 1:
            add_arrow(s, x + bw - 0.02, y + 0.18, gap + 0.04, 0.2, color=COL_ARROW_GRAY)

# ====================================================================
# SLIDE 20: 10-12 Gotham / FDE / Bootcamp
# ====================================================================
def slide_20_gotham_fde():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 12, "17", "10-12  Gotham · FDE · Bootcamp")

    # 3 columns
    cols = [
        ("Gotham", "任务运营系统", COL_RED, [
            "情报 + 传感器数据整合",
            "地理空间 + 目标 + 人员",
            "任务规划 + 资产调度",
            "传感器控制（无人机/卫星）",
            "人在回路控制",
        ], "国防任务操作系统，而非只是数据分析工具"),
        ("FDE", "前向部署工程", COL_GREEN, [
            "业务理解 + 数据建模",
            "软件开发 + 产品配置",
            "系统集成 + 场景验证",
            "客户现场即产品研发",
            "\"人类版本的反向传播\"",
        ], "不是实施工程师，而是产品研发的一部分"),
        ("Bootcamp", "验证 + 交付 + 销售", COL_PURPLE, [
            "使用客户真实数据",
            "客户人员直接参与",
            "数小时/数天形成工作流",
            "不做通用产品演示",
            "2024年超 500 场",
        ], "同时承担产品验证、交付和销售"),
    ]
    for i, (name, role, c, items, footnote) in enumerate(cols):
        x = 0.3 + i * 4.3
        w = 4.1
        h = 5.2
        add_rect(s, x, 1.6, w, h, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_plain_rect(s, x, 1.6, w, 0.08, fill_color=c)
        add_text(s, x + 0.2, 1.75, w - 0.4, 0.4, name, size=18, color=c, bold=True)
        add_text(s, x + 0.2, 2.2, w - 0.4, 0.3, role, size=11, color=COL_TEXT_DARK)
        for j, item in enumerate(items):
            y = 2.7 + j * 0.5
            add_circle(s, x + 0.25, y + 0.05, 0.1, fill_color=c)
            add_text(s, x + 0.45, y, w - 0.6, 0.4, item, size=9, color=COL_TEXT_DARK)
        # Footnote
        add_rect(s, x + 0.1, 5.7, w - 0.2, 1.0, fill_color=COL_LIGHT_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_text(s, x + 0.2, 5.78, w - 0.4, 0.8, footnote, size=8, color=COL_TEXT_GRAY)

# ====================================================================
# SLIDE 21: FDE detail (was image)
# ====================================================================
def slide_21_fde_detail():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 12, "17", "11  FDE：前向部署工程详解")

    # Central: FDE role
    add_circle(s, 5.3, 2.8, 2.3, fill_color=COL_GREEN, line_color=COL_GREEN, line_width=2.0)
    add_text(s, 5.3, 3.1, 2.3, 0.4, "FDE", size=22, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 5.3, 3.6, 2.3, 0.3, "Forward Deployed", size=11, color=COL_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, 5.3, 3.9, 2.3, 0.3, "Engineer", size=11, color=COL_LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # 7 responsibilities around FDE
    resp_items = [
        ("业务理解", "深入客户场景", COL_ICON_BLUE, 1.0, 1.2),
        ("数据建模", "设计 Ontology", COL_GREEN, 8.0, 1.2),
        ("软件开发", "Workshop + Functions", COL_ORANGE, 0.3, 3.0),
        ("产品配置", "定制化部署", COL_PURPLE, 9.5, 3.0),
        ("系统集成", "连接外部系统", COL_TEAL, 0.3, 5.0),
        ("用户交互", "培训 + 支持", COL_RED, 9.5, 5.0),
        ("场景验证", "端到端测试", COL_SECTION_HEADER, 4.5, 6.0),
    ]
    for title, desc, c, x, y in resp_items:
        add_rect(s, x, y, 2.8, 0.9, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 0.06, 0.9, fill_color=c)
        add_text(s, x + 0.15, y + 0.1, 2.5, 0.3, title, size=10, color=c, bold=True)
        add_text(s, x + 0.15, y + 0.4, 2.5, 0.3, desc, size=8, color=COL_TEXT_GRAY)

    # Bottom: "人类版本的反向传播"
    add_rect(s, 0.3, 6.2, 12.7, 0.8, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.5)
    add_text(s, 0.5, 6.25, 12.3, 0.3, '"人类版本的反向传播"', size=12, color=COL_GREEN, bold=True)
    add_text(s, 0.5, 6.55, 12.3, 0.35,
             "FDE 在客户现场发现问题 → 回到产品团队 → 改进平台 → 新能力回流到所有客户 — 形成产品研发的闭环。",
             size=9, color=COL_TEXT_DARK)

# ====================================================================
# SLIDE 22: 13 商业模式
# ====================================================================
def slide_22_business():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 13, "17", "13  商业模式：不是纯 SaaS")

    # Left: Revenue
    add_text(s, 0.3, 1.6, 6.0, 0.3, "收入构成", size=13, color=COL_SECTION_HEADER, bold=True)
    revenue = [
        ("Palantir Cloud 订阅", "SaaS 模式", COL_ICON_BLUE),
        ("本地/私有部署", "软件订阅", COL_GREEN),
        ("持续运营和维护", "服务收入", COL_ORANGE),
        ("专业服务", "培训/配置/建模", COL_PURPLE),
    ]
    for i, (title, desc, c) in enumerate(revenue):
        y = 2.0 + i * 0.6
        add_rect(s, 0.3, y, 5.8, 0.5, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.5)
        add_plain_rect(s, 0.3, y, 0.06, 0.5, fill_color=c)
        add_text(s, 0.5, y + 0.05, 3.5, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 0.5, y + 0.28, 3.5, 0.2, desc, size=8, color=COL_TEXT_GRAY)

    add_rect(s, 0.3, 4.5, 5.8, 0.5, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=0.5)
    add_text(s, 0.5, 4.55, 5.4, 0.35,
             "= 平台订阅 + 持续运维 + 高价值工程服务",
             size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)

    # Land and Expand
    add_text(s, 0.3, 5.3, 6.0, 0.3, "Land and Expand", size=13, color=COL_SECTION_HEADER, bold=True)
    expand_flow = ["一个用例", "→ 一个部门", "→ 更多数据源", "→ 更多对象", "→ 更多用户", "→ 更大合同"]
    for i, item in enumerate(expand_flow):
        y = 5.7 + (i // 3) * 0.5
        x = 0.3 + (i % 3) * 2.0
        add_rect(s, x, y, 1.85, 0.4, fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.5)
        add_text(s, x, y + 0.07, 1.85, 0.25, item, size=8, color=COL_TEXT_DARK, align=PP_ALIGN.CENTER)

    # Key metric
    add_metric_box(s, 0.3, 6.8, 2.8, 0.5, "$9,390万", "2025 Top20 客户均收入", COL_RED)

    # Right: Government + Commercial
    add_text(s, 6.6, 1.6, 6.4, 0.3, "政府 + 商业双轮（2025 年 $44.75 亿）", size=13, color=COL_SECTION_HEADER, bold=True)

    # Government
    add_rect(s, 6.6, 2.0, 3.1, 3.5, fill_color=COL_RED_BG, line_color=COL_RED, line_width=1.0)
    add_text(s, 6.8, 2.1, 2.7, 0.3, "政府收入", size=12, color=COL_RED, bold=True)
    add_text(s, 6.8, 2.5, 2.7, 0.4, "$24.02 亿", size=22, color=COL_RED, bold=True)
    add_text(s, 6.8, 3.0, 2.7, 0.25, "占比 54%", size=10, color=COL_TEXT_DARK)
    add_text(s, 6.8, 3.4, 2.7, 0.2, "特点：", size=9, color=COL_RED, bold=True)
    govt_items = ["高安全要求", "任务关键", "长周期合同", "深入集成"]
    for i, item in enumerate(govt_items):
        add_text(s, 6.9, 3.7 + i * 0.32, 2.6, 0.3, f"• {item}", size=9, color=COL_TEXT_GRAY)

    # Commercial
    add_rect(s, 9.9, 2.0, 3.1, 3.5, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=1.0)
    add_text(s, 10.1, 2.1, 2.7, 0.3, "商业收入", size=12, color=COL_GREEN, bold=True)
    add_text(s, 10.1, 2.5, 2.7, 0.4, "$20.73 亿", size=22, color=COL_GREEN, bold=True)
    add_text(s, 10.1, 3.0, 2.7, 0.25, "占比 46%", size=10, color=COL_TEXT_DARK)
    add_text(s, 10.1, 3.4, 2.7, 0.2, "特点：", size=9, color=COL_GREEN, bold=True)
    comm_items = ["行业场景丰富", "增长空间大", "Land & Expand", "跨行业复制"]
    for i, item in enumerate(comm_items):
        add_text(s, 10.2, 3.7 + i * 0.32, 2.6, 0.3, f"• {item}", size=9, color=COL_TEXT_GRAY)

    # Shared platform
    add_rect(s, 6.6, 5.7, 6.4, 1.1, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=0.75)
    add_text(s, 6.8, 5.78, 6.0, 0.3, "共用平台", size=11, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 6.8, 6.1, 6.0, 0.35,
             "政府线和商业线共用同一套 Foundry + Ontology + AIP 底座，差异在于部署环境和安全等级。",
             size=9, color=COL_TEXT_DARK)

    add_text(s, 6.6, 7.0, 6.4, 0.25,
             "同比增速：商业 +31%，政府 +17%（2025 年）",
             size=8, color=COL_TEXT_GRAY)

# ====================================================================
# SLIDE 23: Revenue breakdown (was image)
# ====================================================================
def slide_23_revenue():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 14, "17", "14  经营数据")

    # Key metrics row
    metrics = [
        ("$44.75亿", "2025 总收入", COL_RED),
        ("+28%", "YoY 增长", COL_GREEN),
        ("$9,390万", "Top20 客户均收入", COL_ORANGE),
        ("954+", "总客户数", COL_SECTION_HEADER),
        ("82%", "2025 毛利率", COL_ICON_BLUE),
        ("87%", "2026 Q1 毛利率", COL_PURPLE),
    ]
    for i, (val, label, c) in enumerate(metrics):
        x = 0.3 + i * 2.13
        add_rect(s, x, 1.6, 2.0, 1.2, fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_text(s, x + 0.05, 1.7, 1.9, 0.45, val, size=18, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, 2.2, 1.9, 0.3, label, size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Revenue trend
    add_text(s, 0.3, 3.1, 6.0, 0.3, "收入趋势", size=13, color=COL_SECTION_HEADER, bold=True)
    years = [("2021", 1.54, COL_TEXT_LIGHT), ("2022", 1.91, COL_TEXT_LIGHT),
             ("2023", 2.23, COL_TEXT_LIGHT), ("2024", 2.87, COL_ORANGE),
             ("2025", 4.48, COL_RED)]
    max_val = 4.48
    chart_x = 0.5
    chart_y = 5.8
    chart_w = 5.5
    bar_w = 0.8
    bar_gap = 0.3
    for i, (year, val, c) in enumerate(years):
        x = chart_x + i * (bar_w + bar_gap)
        h = (val / max_val) * 2.2
        y = chart_y - h
        add_plain_rect(s, x, y, bar_w, h, fill_color=c)
        add_text(s, x - 0.1, y - 0.3, bar_w + 0.2, 0.25, f"${val}B", size=8, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, chart_y + 0.05, bar_w, 0.25, year, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Axis
    add_line(s, chart_x - 0.1, chart_y + 0.3, chart_x + 5.5, chart_y + 0.3, color=COL_BORDER_GRAY, width=1.0)

    # Right: Customer breakdown
    add_text(s, 6.8, 3.1, 6.2, 0.3, "客户构成", size=13, color=COL_SECTION_HEADER, bold=True)

    add_rect(s, 6.8, 3.5, 6.2, 1.3, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.5)
    add_text(s, 7.0, 3.6, 5.8, 0.25, "政府客户", size=11, color=COL_RED, bold=True)
    add_text(s, 7.0, 3.9, 5.8, 0.25, "$24.02 亿（54%）", size=10, color=COL_TEXT_DARK)
    add_text(s, 7.0, 4.2, 5.8, 0.25, "国防 · 情报 · 民用政府", size=9, color=COL_TEXT_GRAY)
    add_text(s, 7.0, 4.45, 5.8, 0.25, "特点：高 ARPU · 长合同周期 · 深度集成", size=8, color=COL_TEXT_GRAY)

    add_rect(s, 6.8, 4.9, 6.2, 1.3, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.5)
    add_text(s, 7.0, 5.0, 5.8, 0.25, "商业客户", size=11, color=COL_GREEN, bold=True)
    add_text(s, 7.0, 5.3, 5.8, 0.25, "$20.73 亿（46%）", size=10, color=COL_TEXT_DARK)
    add_text(s, 7.0, 5.6, 5.8, 0.25, "制造 · 金融 · 医疗 · 能源 · 零售", size=9, color=COL_TEXT_GRAY)
    add_text(s, 7.0, 5.85, 5.8, 0.25, "特点：快速扩张 · 行业复制 · Land & Expand", size=8, color=COL_TEXT_GRAY)

    # TCV
    add_rect(s, 6.8, 6.3, 6.2, 0.8, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=0.5)
    add_text(s, 7.0, 6.35, 5.8, 0.3, "2025 年新签 TCV（总合同价值）：", size=10, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 7.0, 6.65, 5.8, 0.3, "总计 $39.7 亿，商业 TCV $31.4 亿（79%）", size=9, color=COL_TEXT_DARK)

# ====================================================================
# SLIDE 24: Revenue detail (was image)
# ====================================================================
def slide_24_revenue_detail():
    s = prs.slides.add_slide(prs.slides_layouts[6] if hasattr(prs, 'slides_layouts') else prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 14, "17", "14  经营数据：增长指标")

    # Left: Growth metrics
    add_text(s, 0.3, 1.6, 6.0, 0.3, "关键增长指标", size=13, color=COL_SECTION_HEADER, bold=True)
    growth = [
        ("收入增长", "+28% YoY", "$3.50B → $4.48B", COL_GREEN),
        ("商业收入增长", "+31% YoY", "$1.58B → $2.07B", COL_GREEN),
        ("政府收入增长", "+17% YoY", "$2.05B → $2.40B", COL_ORANGE),
        ("客户数增长", "+43% YoY", "668 → 954+", COL_ICON_BLUE),
        ("Top20 客户均收入", "$9,390万", "同比 +16%", COL_RED),
        ("商业 TCV", "$31.4亿", "占新签 79%", COL_PURPLE),
    ]
    for i, (title, value, detail, c) in enumerate(growth):
        y = 2.0 + i * 0.75
        add_rect(s, 0.3, y, 6.0, 0.65, fill_color=COL_BG_WHITE, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_plain_rect(s, 0.3, y, 0.06, 0.65, fill_color=c)
        add_text(s, 0.5, y + 0.03, 2.5, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 0.5, y + 0.3, 2.5, 0.25, detail, size=8, color=COL_TEXT_GRAY)
        add_text(s, 3.2, y + 0.1, 2.7, 0.4, value, size=14, color=c, bold=True, align=PP_ALIGN.RIGHT)

    # Right: Margin trend
    add_text(s, 6.6, 1.6, 6.4, 0.3, "毛利率提升趋势", size=13, color=COL_SECTION_HEADER, bold=True)
    margins = [("2023", 80, COL_TEXT_LIGHT), ("2024", 81, COL_TEXT_LIGHT),
               ("2025", 82, COL_ORANGE), ("2026 Q1", 87, COL_RED)]
    chart_x = 7.0
    chart_y = 5.0
    bar_w = 1.2
    bar_gap = 0.3
    max_h = 2.5
    for i, (label, val, c) in enumerate(margins):
        x = chart_x + i * (bar_w + bar_gap)
        h = (val / 100) * max_h
        y = chart_y - h
        add_plain_rect(s, x, y, bar_w, h, fill_color=c)
        add_text(s, x - 0.1, y - 0.3, bar_w + 0.2, 0.25, f"{val}%", size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, chart_y + 0.05, bar_w, 0.25, label, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_line(s, chart_x - 0.2, chart_y + 0.3, chart_x + 5.5, chart_y + 0.3, color=COL_BORDER_GRAY, width=1.0)

    # Key insight
    add_rect(s, 6.6, 5.6, 6.4, 1.4, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 6.8, 5.7, 6.0, 0.3, "关键洞察：", size=11, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 6.8, 6.05, 6.0, 0.4,
             "毛利率从 80% 持续提升到 87%，说明规模效应开始显现——软件订阅收入占比增大，",
             size=9, color=COL_TEXT_DARK)
    add_text(s, 6.8, 6.45, 6.0, 0.35,
             "专业服务成本占比下降，边际成本递减。",
             size=9, color=COL_TEXT_DARK)

    # Bottom summary
    add_rect(s, 0.3, 6.6, 6.0, 0.5, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.5)
    add_text(s, 0.5, 6.7, 5.6, 0.3,
             "收入质量：高毛利 + 高续约 + 高 ARPU = 优质 SaaS+ 模型",
             size=10, color=COL_GREEN, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 25-26: More charts (were images)
# ====================================================================
def slide_25_customer_growth():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 14, "17", "14  经营数据：客户增长")

    # Customer count trend
    add_text(s, 0.3, 1.6, 6.0, 0.3, "客户数量增长", size=13, color=COL_SECTION_HEADER, bold=True)
    cust_years = [("2021", 237, COL_TEXT_LIGHT), ("2022", 367, COL_TEXT_LIGHT),
                  ("2023", 457, COL_TEXT_LIGHT), ("2024", 668, COL_ORANGE),
                  ("2025", 954, COL_RED)]
    chart_x = 0.8
    chart_y = 5.0
    bar_w = 0.9
    bar_gap = 0.25
    max_val = 954
    max_h = 2.5
    for i, (year, val, c) in enumerate(cust_years):
        x = chart_x + i * (bar_w + bar_gap)
        h = (val / max_val) * max_h
        y = chart_y - h
        add_plain_rect(s, x, y, bar_w, h, fill_color=c)
        add_text(s, x - 0.1, y - 0.3, bar_w + 0.2, 0.25, str(val), size=9, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, chart_y + 0.05, bar_w, 0.25, year, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_line(s, chart_x - 0.2, chart_y + 0.3, chart_x + 5.5, chart_y + 0.3, color=COL_BORDER_GRAY, width=1.0)
    add_text(s, 0.3, 5.4, 5.5, 0.25, "+43% YoY（2024→2025）", size=9, color=COL_GREEN, bold=True)

    # Right: ARPU
    add_text(s, 6.6, 1.6, 6.4, 0.3, "Top20 客户平均收入", size=13, color=COL_SECTION_HEADER, bold=True)
    arpu_years = [("2023", 5400, COL_TEXT_LIGHT), ("2024", 8100, COL_TEXT_LIGHT), ("2025", 9390, COL_RED)]
    chart_x2 = 7.5
    chart_y2 = 5.0
    bar_w2 = 1.2
    bar_gap2 = 0.4
    max_val2 = 9390
    for i, (year, val, c) in enumerate(arpu_years):
        x = chart_x2 + i * (bar_w2 + bar_gap2)
        h = (val / max_val2) * max_h
        y = chart_y2 - h
        add_plain_rect(s, x, y, bar_w2, h, fill_color=c)
        add_text(s, x - 0.1, y - 0.3, bar_w2 + 0.2, 0.25, f"${val/1000:.1f}M", size=9, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, chart_y2 + 0.05, bar_w2, 0.25, year, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_line(s, chart_x2 - 0.2, chart_y2 + 0.3, chart_x2 + 5.0, chart_y2 + 0.3, color=COL_BORDER_GRAY, width=1.0)
    add_text(s, 6.6, 5.4, 5.5, 0.25, "+16% YoY（2024→2025）", size=9, color=COL_GREEN, bold=True)

    # Bottom insight
    add_rect(s, 0.3, 5.9, 12.7, 1.0, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 6.0, 12.3, 0.3, "增长质量：", size=11, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 0.5, 6.35, 12.3, 0.4,
             "客户数增长 +43%（量增）同时 Top20 ARPU 增长 +16%（价增）— 既在扩张新客户，也在深化现有客户。",
             size=9, color=COL_TEXT_DARK)

def slide_26_geo():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 14, "17", "14  经营数据：区域分布")

    # US vs International
    add_text(s, 0.3, 1.6, 12.7, 0.3, "收入地理分布", size=13, color=COL_SECTION_HEADER, bold=True)

    regions = [
        ("美国", "$37.2亿", "83%", COL_SECTION_HEADER, 10.5),
        ("国际", "$7.6亿", "17%", COL_GREEN, 2.1),
    ]
    # Stacked bar
    total_bar_w = 12.7
    x = 0.3
    for name, val, pct, c, w in regions:
        add_plain_rect(s, x, 2.1, w, 0.8, fill_color=c)
        add_text(s, x, 2.3, w, 0.4, f"{name}\n{val} · {pct}", size=11, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w

    # Detail breakdown
    add_text(s, 0.3, 3.3, 12.7, 0.3, "美国市场细分", size=12, color=COL_SECTION_HEADER, bold=True)
    us_items = [
        ("美国政府", "$17.2亿", "46%", COL_RED),
        ("美国商业", "$20.0亿", "54%", COL_GREEN),
    ]
    x = 0.3
    total_w = 12.7
    for name, val, pct, c in us_items:
        w = total_w * (int(pct.strip('%')) / 100)
        add_plain_rect(s, x, 3.7, w, 0.6, fill_color=c)
        add_text(s, x, 3.8, w, 0.4, f"{name}: {val} ({pct})", size=10, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w

    # International
    add_text(s, 0.3, 4.6, 12.7, 0.3, "国际市场细分", size=12, color=COL_SECTION_HEADER, bold=True)
    intl = [
        ("英国", "$3.6亿", COL_ICON_BLUE),
        ("欧洲其他", "$2.1亿", COL_GREEN),
        ("亚太", "$1.2亿", COL_ORANGE),
        ("其他", "$0.7亿", COL_PURPLE),
    ]
    intl_total = 7.6
    x = 0.3
    for name, val, c in intl:
        v = float(val.replace('$','').replace('亿',''))
        w = total_w * (v / intl_total) * 0.6
        add_plain_rect(s, x, 5.0, w, 0.6, fill_color=c)
        add_text(s, x, 5.1, w, 0.4, f"{name}\n{val}", size=8, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w

    # Key insight
    add_rect(s, 0.3, 5.9, 12.7, 1.0, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 6.0, 12.3, 0.3, "地理洞察：", size=11, color=COL_SECTION_HEADER, bold=True)
    add_text(s, 0.5, 6.35, 12.3, 0.4,
             "美国市场仍占主导（83%），但国际市场增速更快（+45% YoY）。英国是最大的国际市场。",
             size=9, color=COL_TEXT_DARK)

# ====================================================================
# SLIDE 27: 15 复合飞轮壁垒
# ====================================================================
def slide_27_flywheel():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15  真正壁垒：一套相互增强的系统")

    # 7 pillars
    pillars = [
        ("数据 + 计算", "连接复杂异构\n高敏感数据", COL_SECTION_HEADER),
        ("Ontology", "统一业务\n运行模型", COL_ICON_BLUE),
        ("Action + 反馈", "分析结果\n进入业务执行", COL_GREEN),
        ("AIP", "大模型进入\n已有业务模型", COL_PURPLE),
        ("Apollo + Rubix", "多云/边缘部署\n持续升级", COL_ORANGE),
        ("FDE", "客户问题\n转化为产品能力", COL_TEAL),
        ("Bootcamp", "降低初次\n验证成本", COL_RED),
        ("账户扩张", "一次性问题\n→ 长期平台收入", COL_GOLD),
    ]
    for i, (title, desc, c) in enumerate(pillars):
        col = i % 4
        row = i // 4
        x = 0.3 + col * 3.2
        y = 1.7 + row * 1.6
        add_rect(s, x, y, 3.0, 1.4, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_plain_rect(s, x, y, 3.0, 0.06, fill_color=c)
        add_text(s, x + 0.1, y + 0.12, 2.8, 0.35, title, size=12, color=c, bold=True)
        add_text(s, x + 0.1, y + 0.5, 2.8, 0.8, desc, size=9, color=COL_TEXT_GRAY)

    # Flywheel description
    add_rect(s, 0.3, 5.1, 12.7, 1.8, fill_color=COL_LIGHTER_BLUE, line_color=COL_SECTION_HEADER, line_width=1.0)
    add_text(s, 0.5, 5.2, 12.3, 0.3, "复合飞轮：", size=13, color=COL_SECTION_HEADER, bold=True)
    flywheel = [
        "更多复杂客户问题",
        "→ 更多现场工程经验",
        "→ 更强的平台能力",
        "→ 更快的用例交付",
        "→ 更多客户和场景",
        "→ 更多 Ontology/Action/决策数据",
        "→ 更高平台价值",
        "→ 更深客户扩张",
    ]
    for i, item in enumerate(flywheel):
        col = i % 4
        row = i // 4
        x = 0.5 + col * 3.1
        y = 5.6 + row * 0.55
        add_text(s, x, y, 3.0, 0.45, item, size=9, color=COL_SECTION_HEADER if row == 0 else COL_TEXT_DARK,
                 bold=(row == 0))

    # Conclusion
    add_rect(s, 0.3, 7.0, 12.7, 0.35, fill_color=COL_SECTION_HEADER)
    add_text(s, 0.5, 7.03, 12.3, 0.3,
             "竞争优势不是某个模块全面领先，而是能把多个模块组合成一套完整运营系统。",
             size=10, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 28: 16 约束
# ====================================================================
def slide_28_constraints():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 16, "17", "16  也存在明显约束")

    constraints = [
        ("实施复杂度高", COL_RED, [
            "平台复杂，实施周期长",
            "需要针对客户独特环境配置",
            "需要培训和持续技术人员服务",
        ]),
        ("销售成本高", COL_ORANGE, [
            "大型项目安装成本高",
            "失败风险高",
            "评估周期可能超过一年",
        ]),
        ("依赖工程服务", COL_PURPLE, [
            "成本含现场人员/专业服务",
            "部分客户需长期参与",
            "分包商/云资源/部署维护",
        ]),
        ("面临多类竞争", COL_TEAL, [
            "大型软件公司",
            "政府承包商",
            "系统集成商/新兴公司/自研",
        ]),
    ]
    for i, (title, c, items) in enumerate(constraints):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.4
        y = 1.7 + row * 2.5
        w = 6.2
        h = 2.3
        add_rect(s, x, y, w, h, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 0.08, h, fill_color=c)
        add_text(s, x + 0.25, y + 0.12, 5.5, 0.35, title, size=13, color=c, bold=True)
        for j, item in enumerate(items):
            add_circle(s, x + 0.3, y + 0.65 + j * 0.4 + 0.03, 0.08, fill_color=c)
            add_text(s, x + 0.5, y + 0.6 + j * 0.4, 5.3, 0.35, item, size=9, color=COL_TEXT_GRAY)

    # Bottom: mitigant
    add_rect(s, 0.3, 6.7, 12.7, 0.5, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=0.5)
    add_text(s, 0.5, 6.78, 12.3, 0.35,
             "竞争优势不是某个模块全面领先，而是能把多个模块组合成一套完整运营系统。",
             size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 29-32: Additional chart pages (were images)
# ====================================================================
def slide_29_rnd():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15  复合飞轮壁垒：研发投入")

    # R&D metrics
    add_text(s, 0.3, 1.6, 6.0, 0.3, "研发投入", size=13, color=COL_SECTION_HEADER, bold=True)
    rnd_items = [
        ("2025 R&D 支出", "$6.74亿", "占收入 15%", COL_RED),
        ("FDE 团队规模", "~1,600人", "占员工 36%", COL_GREEN),
        ("专利数量", "180+", "核心技术保护", COL_ICON_BLUE),
        ("平台迭代频率", "每日数千次", "Apollo 部署", COL_ORANGE),
    ]
    for i, (title, val, desc, c) in enumerate(rnd_items):
        y = 2.0 + i * 0.9
        add_rect(s, 0.3, y, 6.0, 0.8, fill_color=COL_BG_WHITE, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_plain_rect(s, 0.3, y, 0.06, 0.8, fill_color=c)
        add_text(s, 0.5, y + 0.05, 2.5, 0.3, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 0.5, y + 0.38, 2.5, 0.25, desc, size=8, color=COL_TEXT_GRAY)
        add_text(s, 3.2, y + 0.15, 2.8, 0.45, val, size=16, color=c, bold=True, align=PP_ALIGN.RIGHT)

    # Right: Product velocity
    add_text(s, 6.6, 1.6, 6.4, 0.3, "产品迭代速度", size=13, color=COL_SECTION_HEADER, bold=True)
    velocity = [
        ("AIP 发布", "2023", "从 0 到企业级 AI 运行环境", COL_PURPLE),
        ("Ontology 增强", "持续", "对象/Action/函数不断扩展", COL_GREEN),
        ("Apollo 部署", "每日", "数千次零停机升级", COL_ORANGE),
        ("Bootcamp", "2024 500+场", "客户验证→产品迭代", COL_RED),
        ("行业模板", "持续积累", "从客户场景提炼通用能力", COL_TEAL),
    ]
    for i, (title, period, desc, c) in enumerate(velocity):
        y = 2.0 + i * 0.7
        add_rect(s, 6.6, y, 6.4, 0.6, fill_color=COL_BG_WHITE, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_circle(s, 6.8, y + 0.08, 0.18, fill_color=c)
        add_text(s, 6.8, y + 0.12, 0.18, 0.2, str(i+1), size=9, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, 7.15, y + 0.03, 2.0, 0.25, title, size=10, color=COL_TITLE_DARK, bold=True)
        add_text(s, 7.15, y + 0.28, 2.0, 0.22, period, size=8, color=c)
        add_text(s, 9.3, y + 0.1, 3.5, 0.35, desc, size=8, color=COL_TEXT_GRAY)

    # Bottom
    add_rect(s, 0.3, 5.8, 12.7, 1.2, fill_color=COL_LIGHTER_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
    add_text(s, 0.5, 5.9, 12.3, 0.3, "FDE = 人类版本的反向传播：", size=11, color=COL_SECTION_HEADER, bold=True)
    fde_items = [
        "业务理解", "数据建模", "软件开发", "产品配置", "系统集成", "用户交互", "场景验证",
    ]
    for i, item in enumerate(fde_items):
        x = 0.5 + i * 1.75
        add_rect(s, x, 6.3, 1.6, 0.5, fill_color=COL_BG_WHITE, line_color=COL_GREEN, line_width=0.75)
        add_text(s, x, 6.4, 1.6, 0.3, item, size=9, color=COL_GREEN, bold=True, align=PP_ALIGN.CENTER)

def slide_30_market():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15  复合飞轮壁垒：市场与竞争")

    # TAM
    add_text(s, 0.3, 1.6, 12.7, 0.3, "可寻址市场", size=13, color=COL_SECTION_HEADER, bold=True)
    tam = [
        ("TAM", "$1,199亿", "全球数据+AI+决策平台", COL_RED),
        ("SAM", "$450亿", "大型政企决策与运营", COL_ORANGE),
        ("SOM", "$40亿", "近期可获取份额", COL_GREEN),
    ]
    for i, (name, val, desc, c) in enumerate(tam):
        x = 0.3 + i * 4.2
        add_rect(s, x, 2.0, 4.0, 1.5, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_plain_rect(s, x, y := 2.0, 4.0, 0.06, fill_color=c)
        add_text(s, x + 0.1, 2.15, 3.8, 0.3, name, size=12, color=c, bold=True)
        add_text(s, x + 0.1, 2.55, 3.8, 0.5, val, size=24, color=c, bold=True)
        add_text(s, x + 0.1, 3.1, 3.8, 0.3, desc, size=9, color=COL_TEXT_GRAY)

    # Competition
    add_text(s, 0.3, 3.8, 12.7, 0.3, "竞争格局", size=13, color=COL_SECTION_HEADER, bold=True)
    competitors = [
        ("大型软件公司", "Microsoft / Oracle / SAP", "有平台但缺深度运营闭环", COL_ICON_BLUE),
        ("云厂商", "AWS / GCP / Azure", "有基础设施但缺业务语义层", COL_GREEN),
        ("数据平台", "Snowflake / Databricks", "有数据处理但缺 Ontology + Action", COL_ORANGE),
        ("SI / 咨询", "Accenture / Deloitte", "有行业经验但缺产品平台", COL_PURPLE),
        ("AI 公司", "OpenAI / Anthropic", "有模型但缺企业运营上下文", COL_TEAL),
        ("客户自研", "内部技术团队", "有业务理解但缺平台和规模", COL_RED),
    ]
    for i, (cat, examples, gap, c) in enumerate(competitors):
        col = i % 3
        row = i // 3
        x = 0.3 + col * 4.2
        y = 4.2 + row * 1.4
        add_rect(s, x, y, 4.0, 1.25, fill_color=COL_BG_WHITE, line_color=c, line_width=0.75)
        add_plain_rect(s, x, y, 0.06, 1.25, fill_color=c)
        add_text(s, x + 0.2, y + 0.08, 3.5, 0.25, cat, size=10, color=c, bold=True)
        add_text(s, x + 0.2, y + 0.35, 3.5, 0.25, examples, size=9, color=COL_TEXT_DARK)
        add_text(s, x + 0.2, y + 0.65, 3.5, 0.35, gap, size=8, color=COL_TEXT_GRAY)

    # Bottom
    add_rect(s, 0.3, 7.0, 12.7, 0.35, fill_color=COL_SECTION_HEADER)
    add_text(s, 0.5, 7.03, 12.3, 0.3,
             "Palantir 的差异化：不是单点最优，而是「数据→Ontology→Action→AI→部署→反馈」的完整闭环。",
             size=9, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)

def slide_31_ecosystem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15  复合飞轮壁垒：生态系统")

    # Partner ecosystem
    add_text(s, 0.3, 1.6, 6.0, 0.3, "合作伙伴生态", size=13, color=COL_SECTION_HEADER, bold=True)
    partners = [
        ("云伙伴", "AWS · GCP · Azure", "多云部署支持", COL_ICON_BLUE),
        ("数据伙伴", "Snowflake · Databricks", "数据湖集成", COL_GREEN),
        ("SI 伙伴", "Accenture · Deloitte · IBM", "规模化交付", COL_ORANGE),
        ("行业伙伴", "SAP · Salesforce · Oracle", "企业系统连接", COL_PURPLE),
        ("AI 伙伴", "OpenAI · Anthropic · Mistral", "模型供给", COL_TEAL),
        ("硬件伙伴", "NVIDIA · Dell · HPE", "边缘/私有部署", COL_RED),
    ]
    for i, (cat, names, desc, c) in enumerate(partners):
        y = 2.0 + i * 0.75
        add_rect(s, 0.3, y, 6.0, 0.65, fill_color=COL_BG_WHITE, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_plain_rect(s, 0.3, y, 0.06, 0.65, fill_color=c)
        add_text(s, 0.5, y + 0.03, 1.8, 0.25, cat, size=10, color=c, bold=True)
        add_text(s, 0.5, y + 0.3, 1.8, 0.25, desc, size=8, color=COL_TEXT_GRAY)
        add_text(s, 2.5, y + 0.1, 3.5, 0.4, names, size=9, color=COL_TEXT_DARK)

    # Right: Moat analysis
    add_text(s, 6.6, 1.6, 6.4, 0.3, "护城河分析（波特五力+）", size=13, color=COL_SECTION_HEADER, bold=True)
    moats = [
        ("数据网络效应", "更多数据→更好模型→更多客户", COL_RED, 9),
        ("转换成本", "Ontology 深度嵌入→迁移困难", COL_ORANGE, 8),
        ("规模经济", "平台固定成本→边际递减", COL_GREEN, 7),
        ("品牌信任", "政府/军事级安全认证", COL_PURPLE, 8),
        ("技术复杂度", "多组件协同→难以复制", COL_TEAL, 9),
        ("客户锁定", "Land & Expand→深度依赖", COL_ICON_BLUE, 8),
    ]
    for i, (title, desc, c, score) in enumerate(moats):
        y = 2.0 + i * 0.75
        add_rect(s, 6.6, y, 6.4, 0.65, fill_color=COL_BG_WHITE, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_plain_rect(s, 6.6, y, 0.06, 0.65, fill_color=c)
        add_text(s, 6.8, y + 0.03, 2.2, 0.25, title, size=10, color=c, bold=True)
        add_text(s, 6.8, y + 0.3, 3.5, 0.25, desc, size=8, color=COL_TEXT_GRAY)
        # Score bar
        bar_x = 10.5
        bar_w = (score / 10) * 2.3
        add_plain_rect(s, bar_x, y + 0.2, 2.3, 0.25, fill_color=COL_BORDER_LIGHT)
        add_plain_rect(s, bar_x, y + 0.2, bar_w, 0.25, fill_color=c)
        add_text(s, bar_x + 2.35, y + 0.17, 0.5, 0.3, f"{score}/10", size=8, color=c, bold=True)

def slide_32_risks():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 16, "17", "16  约束与挑战：深入分析")

    # Risk matrix
    add_text(s, 0.3, 1.6, 12.7, 0.3, "风险评估矩阵", size=13, color=COL_SECTION_HEADER, bold=True)

    risks = [
        ("实施复杂度", "高", COL_RED, [
            "平台组件多，学习曲线陡峭",
            "需要深度业务理解 + 技术能力",
            "FDE 模式导致人力成本高",
        ]),
        ("销售周期", "长", COL_ORANGE, [
            "大型项目评估周期 6-18 个月",
            "POC → Bootcamp → 合同 → 实施",
            "政府项目受预算周期影响",
        ]),
        ("盈利路径", "渐进", COL_PURPLE, [
            "前 1-2 年投入大于收入",
            "需要客户扩展到一定规模后盈利",
            "毛利率改善需要时间",
        ]),
        ("技术债务", "可控", COL_TEAL, [
            "历史组件（Gotham）与新平台并存",
            "多版本兼容性维护成本",
            "需要持续重构和升级",
        ]),
    ]
    for i, (title, level, c, items) in enumerate(risks):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.4
        y = 2.0 + row * 2.5
        w = 6.2
        h = 2.3
        add_rect(s, x, y, w, h, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, w, 0.06, fill_color=c)
        add_text(s, x + 0.25, y + 0.12, 3.5, 0.3, title, size=12, color=c, bold=True)
        # Risk level badge
        add_rect(s, x + w - 1.2, y + 0.12, 0.9, 0.3, fill_color=c)
        add_text(s, x + w - 1.2, y + 0.15, 0.9, 0.25, level, size=9, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_circle(s, x + 0.3, y + 0.6 + j * 0.4 + 0.03, 0.06, fill_color=c)
            add_text(s, x + 0.45, y + 0.55 + j * 0.4, w - 0.7, 0.35, item, size=9, color=COL_TEXT_GRAY)

    # Bottom: Mitigation
    add_rect(s, 0.3, 7.0, 12.7, 0.35, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.5)
    add_text(s, 0.5, 7.03, 12.3, 0.3,
             "缓解策略：Bootcamp 降低验证成本 · AIP 加速用例交付 · Apollo 降低运维复杂度 · 行业模板复用经验。",
             size=9, color=COL_GREEN, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 33-35: Moat details (were images)
# ====================================================================
def slide_33_moat_data():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15.1  数据和计算能力壁垒")

    # Key capabilities
    add_text(s, 0.3, 1.6, 12.7, 0.3, "核心数据能力", size=13, color=COL_SECTION_HEADER, bold=True)
    caps = [
        ("多源连接", "190+ 连接器\nERP/CRM/IoT/传感器", COL_ICON_BLUE),
        ("多模态处理", "结构化/非结构化\n流/地理空间/文档", COL_GREEN),
        ("数据治理", "字段级权限\n血缘追踪/审计", COL_ORANGE),
        ("高安全部署", "多云/私有/边缘\nFedRAMP/IL5", COL_RED),
        ("实时处理", "流式数据\n亚秒级延迟", COL_PURPLE),
        ("弹性计算", "自动扩缩容\nGPU/CPU 混合", COL_TEAL),
    ]
    for i, (title, desc, c) in enumerate(caps):
        col = i % 3
        row = i // 3
        x = 0.3 + col * 4.2
        y = 2.0 + row * 1.5
        add_rect(s, x, y, 4.0, 1.35, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 4.0, 0.06, fill_color=c)
        add_text(s, x + 0.1, y + 0.15, 3.8, 0.3, title, size=11, color=c, bold=True)
        add_text(s, x + 0.1, y + 0.55, 3.8, 0.7, desc, size=9, color=COL_TEXT_GRAY)

    # Why hard to replicate
    add_rect(s, 0.3, 5.2, 12.7, 1.8, fill_color=COL_RED_BG, line_color=COL_RED, line_width=0.5)
    add_text(s, 0.5, 5.3, 12.3, 0.3, "为什么难以复制：", size=11, color=COL_RED, bold=True)
    reasons = [
        "连接器库需要 10+ 年积累（190+ 种，含冷门工业系统）",
        "高安全认证（FedRAMP High / IL5 / DoD IL6）需要数年",
        "多模态数据处理需要深厚的工程能力",
        "在敏感环境（离线/边缘）中运行需要专门的架构设计",
    ]
    for i, r in enumerate(reasons):
        add_text(s, 0.6, 5.7 + i * 0.3, 12.0, 0.25, f"▸ {r}", size=9, color=COL_TEXT_DARK)

def slide_34_moat_ontology():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15.2  Ontology 壁垒")

    # Ontology as moat
    add_text(s, 0.3, 1.6, 12.7, 0.3, "Ontology 为什么是壁垒", size=13, color=COL_SECTION_HEADER, bold=True)

    points = [
        ("统一业务运行模型", "将数据、业务语义、逻辑、Action、权限\n组织成一套模型 — 竞品大多只有数据层", COL_ICON_BLUE),
        ("深度嵌入客户流程", "一旦客户在 Ontology 上构建了\n数百个对象和 Action，迁移成本极高", COL_GREEN),
        ("非静态知识图谱", "是运行时——支持事务、实时修改、\nAction 写回，远超传统 KG 产品", COL_ORANGE),
        ("与 AI 的天然结合", "AIP 直接在 Ontology 上运行，\nAI 拥有完整业务上下文", COL_PURPLE),
        ("行业模板积累", "每个客户场景都在丰富平台的\n行业 Ontology 模板库", COL_TEAL),
        ("Global Branching", "Ontology 修改可分支、测试、\n审核 — 工程化管理业务变化", COL_RED),
    ]
    for i, (title, desc, c) in enumerate(points):
        col = i % 3
        row = i // 3
        x = 0.3 + col * 4.2
        y = 2.0 + row * 2.3
        add_rect(s, x, y, 4.0, 2.1, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 0.06, 2.1, fill_color=c)
        add_text(s, x + 0.2, y + 0.12, 3.6, 0.35, title, size=11, color=c, bold=True)
        add_text(s, x + 0.2, y + 0.55, 3.6, 1.4, desc, size=9, color=COL_TEXT_GRAY)

    # Bottom
    add_rect(s, 0.3, 6.8, 12.7, 0.4, fill_color=COL_LIGHT_BLUE, line_color=COL_SECTION_HEADER, line_width=0.5)
    add_text(s, 0.5, 6.88, 12.3, 0.3,
             "Ontology 是 Palantir 最核心的差异化 — 它不是数据库、不是知识图谱，而是企业业务的运行时。",
             size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)

def slide_35_moat_fde():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 15, "17", "15.3  FDE 与反馈飞轮壁垒")

    # FDE cycle
    add_text(s, 0.3, 1.6, 12.7, 0.3, "FDE 驱动的产品飞轮", size=13, color=COL_SECTION_HEADER, bold=True)

    cycle = [
        ("客户现场", "FDE 部署到客户\n发现真实问题", COL_ICON_BLUE),
        ("问题抽象", "将个案抽象为\n通用模式", COL_GREEN),
        ("产品迭代", "平台增加新能力\n解决一类问题", COL_ORANGE),
        ("能力回流", "新能力可供\n所有客户使用", COL_PURPLE),
        ("更多客户", "更强的平台\n吸引新客户", COL_TEAL),
        ("更多场景", "更多场景\n更多 FDE 部署", COL_RED),
    ]
    bw = 1.95
    gap = 0.1
    start_x = 0.3 + (12.7 - bw * 6 - gap * 5) / 2
    for i, (title, desc, c) in enumerate(cycle):
        x = start_x + i * (bw + gap)
        y = 2.1
        add_rect(s, x, y, bw, 1.4, fill_color=COL_BG_WHITE, line_color=c, line_width=1.5)
        add_circle(s, x + bw/2 - 0.2, y + 0.15, 0.4, fill_color=c)
        add_text(s, x + bw/2 - 0.2, y + 0.22, 0.4, 0.3, str(i+1), size=16, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.65, bw - 0.1, 0.3, title, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, y + 0.95, bw - 0.1, 0.4, desc, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
        if i < len(cycle) - 1:
            add_arrow(s, x + bw - 0.02, y + 0.55, gap + 0.04, 0.25)

    # FDE profile
    add_text(s, 0.3, 3.8, 6.0, 0.3, "FDE 的多重角色", size=12, color=COL_SECTION_HEADER, bold=True)
    roles = [
        ("业务理解", "深入客户运营场景", COL_ICON_BLUE),
        ("数据建模", "设计 Ontology Schema", COL_GREEN),
        ("软件开发", "Workshop + Functions", COL_ORANGE),
        ("产品配置", "定制部署和集成", COL_PURPLE),
        ("系统集成", "连接 ERP/CRM/IoT", COL_TEAL),
        ("用户交互", "培训和持续支持", COL_RED),
        ("场景验证", "端到端测试验证", COL_SECTION_HEADER),
    ]
    for i, (title, desc, c) in enumerate(roles):
        y = 4.2 + i * 0.4
        add_circle(s, 0.5, y + 0.04, 0.12, fill_color=c)
        add_text(s, 0.75, y, 2.0, 0.3, title, size=9, color=COL_TITLE_DARK, bold=True)
        add_text(s, 2.8, y, 3.2, 0.3, desc, size=8, color=COL_TEXT_GRAY)

    # Right: Metrics
    add_text(s, 6.6, 3.8, 6.4, 0.3, "FDE 规模与效率", size=12, color=COL_SECTION_HEADER, bold=True)
    fde_metrics = [
        ("~1,600", "FDE 团队规模", COL_GREEN),
        ("36%", "占员工比例", COL_ICON_BLUE),
        ("500+", "2024 Bootcamp 场次", COL_ORANGE),
        ("2-5天", "典型 Bootcamp 周期", COL_PURPLE),
        ("70%+", "Bootcamp 转化率", COL_RED),
    ]
    for i, (val, label, c) in enumerate(fde_metrics):
        y = 4.2 + i * 0.5
        add_rect(s, 6.6, y, 6.4, 0.42, fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
        add_text(s, 6.8, y + 0.05, 2.0, 0.3, val, size=14, color=c, bold=True)
        add_text(s, 9.0, y + 0.08, 3.8, 0.3, label, size=9, color=COL_TEXT_GRAY)

    # Bottom quote
    add_rect(s, 0.3, 6.8, 12.7, 0.4, fill_color=COL_GREEN_BG, line_color=COL_GREEN, line_width=0.5)
    add_text(s, 0.5, 6.88, 12.3, 0.3,
             '"人类版本的反向传播" — FDE 是 Palantir 产品研发的核心机制，也是最深护城河之一。',
             size=10, color=COL_GREEN, bold=True, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 36: 17 最终研究结论
# ====================================================================
def slide_36_conclusions():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)
    add_title_bar(s, 17, "17", "17  最终研究结论")

    conclusions = [
        ("一", "Palantir 不是传统数据公司", "数据是起点不是终点 — 最终处理的是数据→决策→行动→改变现实", COL_SECTION_HEADER),
        ("二", "Ontology 是核心但不是全部", "只有与数据/逻辑/Action/AI/安全/部署结合后才具运营价值", COL_ICON_BLUE),
        ("三", "AIP 不是替代而是放大器", "用大模型提高企业系统的自然语言交互/推理/工具调用/自动化", COL_PURPLE),
        ("四", "Apollo 决定能否进入关键环境", "能部署到大型企业/政府/军事/私有环境/边缘现场并持续维护", COL_ORANGE),
        ("五", "FDE 是产品研发方法", "客户现场是产品研发循环的一部分", COL_GREEN),
        ("六", "商业模式建立在账户扩张", "从复杂客户高投入开始，逐渐进入更多部门/数据/运营流程", COL_RED),
    ]
    for i, (num, title, desc, c) in enumerate(conclusions):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.4
        y = 1.6 + row * 1.4
        w = 6.2
        h = 1.25
        add_rect(s, x, y, w, h, fill_color=COL_BG_WHITE, line_color=c, line_width=1.0)
        add_plain_rect(s, x, y, 0.08, h, fill_color=c)
        add_circle(s, x + 0.2, y + 0.3, 0.55, fill_color=c)
        add_text(s, x + 0.2, y + 0.38, 0.55, 0.4, num, size=18, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.9, y + 0.12, 5.0, 0.35, title, size=12, color=c, bold=True)
        add_text(s, x + 0.9, y + 0.5, 5.0, 0.6, desc, size=9, color=COL_TEXT_GRAY)

    # Final conclusion
    add_rect(s, 0.3, 5.9, 12.7, 1.2, fill_color=COL_SECTION_HEADER)
    add_text(s, 0.5, 6.0, 12.3, 0.4, "结论七：Palantir = 企业级决策与运营基础设施", size=16, color=COL_BG_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 0.5, 6.5, 12.3, 0.45,
             "以 Ontology 为业务运行模型，以 Foundry 为数据与运营平台，以 AIP 为 AI 运行环境，\n"
             "以 Apollo 和 Rubix 为软件交付底座，连接数据、决策和现实行动。",
             size=11, color=COL_LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 37: BACK COVER
# ====================================================================
def slide_37_back_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_white(s)

    # Top decorative bar
    add_plain_rect(s, 0, 0, 13.33, 0.08, fill_color=COL_SECTION_HEADER)

    # Main text
    add_text(s, 1.0, 1.8, 11.3, 2.0,
             [("以 Ontology 为业务运行模型\n", 16, COL_TITLE_DARK, True),
              ("以 Foundry 为数据与运营平台\n", 16, COL_TITLE_DARK, True),
              ("以 AIP 为 AI 运行环境\n", 16, COL_TITLE_DARK, True),
              ("以 Apollo 和 Rubix 为软件交付底座\n", 16, COL_TITLE_DARK, True),
              ("连接数据、决策和现实行动的\n", 16, COL_TITLE_DARK, True),
              ("企业级决策与运营基础设施", 18, COL_SECTION_HEADER, True)],
             align=PP_ALIGN.CENTER, line_spacing=1.6)

    # Divider
    add_line(s, 4.0, 4.5, 9.33, 4.5, color=COL_BORDER_BLUE, width=1.0)

    # Source
    add_text(s, 1.0, 4.8, 11.3, 0.5,
             "来源：今日头条「23页PPT讲透Palantir」 · 2026.07",
             size=12, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Based on public materials
    add_text(s, 1.0, 5.5, 11.3, 0.35,
             "基于公开资料整理",
             size=10, color=COL_TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Bottom decorative bar
    add_plain_rect(s, 0, 7.42, 13.33, 0.08, fill_color=COL_BORDER_BLUE)


# ============ Generate All Slides ============
print("Generating slides...")

slide_01_cover()
print("  Slide 1: Cover ✓")

slide_02_infographic()
print("  Slide 2: Infographic ✓")

slide_03_toc()
print("  Slide 3: TOC ✓")

slide_04_understand()
print("  Slide 4: Understand ✓")

slide_05_history()
print("  Slide 5: History ✓")

slide_06_platforms()
print("  Slide 6: Platforms ✓")

slide_07_foundry()
print("  Slide 7: Foundry ✓")

slide_08_ontology()
print("  Slide 8: Ontology ✓")

slide_09_ontology_detail()
print("  Slide 9: Ontology Detail ✓")

slide_10_action_loop()
print("  Slide 10: Action Loop ✓")

slide_11_action_types()
print("  Slide 11: Action Types ✓")

slide_12_action()
print("  Slide 12: Action ✓")

slide_13_aip()
print("  Slide 13: AIP ✓")

slide_14_aip_arch()
print("  Slide 14: AIP Architecture ✓")

slide_15_aip_agent()
print("  Slide 15: AIP Agent ✓")

slide_16_apollo()
print("  Slide 16: Apollo & Rubix ✓")

slide_17_apollo_detail()
print("  Slide 17: Apollo Detail ✓")

slide_18_branching()
print("  Slide 18: Global Branching ✓")

slide_19_branching_detail()
print("  Slide 19: Branching Detail ✓")

slide_20_gotham_fde()
print("  Slide 20: Gotham/FDE/Bootcamp ✓")

slide_21_fde_detail()
print("  Slide 21: FDE Detail ✓")

slide_22_business()
print("  Slide 22: Business Model ✓")

slide_23_revenue()
print("  Slide 23: Revenue ✓")

slide_24_revenue_detail()
print("  Slide 24: Revenue Detail ✓")

slide_25_customer_growth()
print("  Slide 25: Customer Growth ✓")

slide_26_geo()
print("  Slide 26: Geo ✓")

slide_27_flywheel()
print("  Slide 27: Flywheel ✓")

slide_28_constraints()
print("  Slide 28: Constraints ✓")

slide_29_rnd()
print("  Slide 29: R&D ✓")

slide_30_market()
print("  Slide 30: Market ✓")

slide_31_ecosystem()
print("  Slide 31: Ecosystem ✓")

slide_32_risks()
print("  Slide 32: Risks ✓")

slide_33_moat_data()
print("  Slide 33: Moat - Data ✓")

slide_34_moat_ontology()
print("  Slide 34: Moat - Ontology ✓")

slide_35_moat_fde()
print("  Slide 35: Moat - FDE ✓")

slide_36_conclusions()
print("  Slide 36: Conclusions ✓")

slide_37_back_cover()
print("  Slide 37: Back Cover ✓")

prs.save(OUTPUT)
print(f"\nDone! {len(prs.slides)} slides saved to: {OUTPUT}")
print(f"File size: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
