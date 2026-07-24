#!/usr/bin/env python3
"""
Test: Convert slide 2 image to native PPTX components with white background.
Slide 2 was: "Palantir 究竟是一家什么公司" - cover info diagram.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = "/Users/ddt/work/projects/ai_agent/docs/ref/AOS_test_slide2.pptx"

# Colors from original image
COL_TITLE_DARK = RGBColor(0x10, 0x35, 0x5E)     # Dark navy/teal title
COL_SUB_BLUE = RGBColor(0x1F, 0x6F, 0xC1)       # Bright blue subtitle
COL_BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_LIGHT_BLUE = RGBColor(0xE9, 0xF2, 0xFC)     # Light blue background
COL_SECTION_HEADER = RGBColor(0x12, 0x47, 0x9A) # Section header blue
COL_ICON_BLUE = RGBColor(0x1F, 0x6F, 0xC1)
COL_BORDER_BLUE = RGBColor(0xC5, 0xDB, 0xF0)
COL_MEDIUM_BLUE = RGBColor(0x4D, 0x88, 0xD6)
COL_TEXT_DARK = RGBColor(0x1F, 0x2D, 0x3D)
COL_TEXT_GRAY = RGBColor(0x5A, 0x6B, 0x7E)
COL_GREEN = RGBColor(0x2D, 0xC2, 0x76)
COL_RED_ACCENT = RGBColor(0xE5, 0x4A, 0x4A)
COL_LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
COL_BORDER_GRAY = RGBColor(0xD8, 0xDF, 0xE8)
COL_ARROW_GRAY = RGBColor(0xB0, 0xBD, 0xCC)
COL_VALUE_BG = RGBColor(0xFA, 0xFB, 0xFD)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

def set_bg_white(s):
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COL_BG_WHITE

def add_text(slide, x, y, w, h, text, size=12, color=None, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Microsoft YaHei"):
    if color is None:
        color = COL_TEXT_DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, list):
        for i, line in enumerate(text):
            if i > 0:
                p = tf.add_paragraph()
                p.alignment = align
            run = p.add_run()
            run.text = line[0]
            run.font.size = Pt(line[1] if len(line) > 1 else size)
            run.font.color.rgb = line[2] if len(line) > 2 else color
            run.font.bold = line[3] if len(line) > 3 else bold
            run.font.name = font_name
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return tb

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    # Adjust corner radius via XML
    sp = shape._element
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        # set moderate corner radius
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 25000')

    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_arrow(slide, x, y, w, h, direction="right", color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        color = COL_ARROW_GRAY
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


set_bg_white(slide)

# ====================================================================
# TOP TITLE BAR
# ====================================================================
# Title
add_text(slide, 0.3, 0.2, 12.7, 0.6,
         "Palantir 究竟是一家什么公司",
         size=28, color=COL_TITLE_DARK, bold=True)

# Subtitle
add_text(slide, 0.3, 0.85, 12.7, 0.4,
         "以数据、业务逻辑和AI为核心，连接决策与现实行动的企业级决策与运营基础设施",
         size=13, color=COL_SUB_BLUE, bold=False)

# Decorative line under title
line_shape = slide.shapes.add_connector(1, Inches(0.3), Inches(1.35),
                                       Inches(13.0), Inches(1.35))
line_shape.line.color.rgb = COL_BORDER_BLUE
line_shape.line.width = Pt(1)

# ====================================================================
# ROW 1: Left intro paragraph + Right "为客户创造的核心价值"
# ====================================================================
# LEFT: Intro paragraph card
intro_bg = add_rect(slide, 0.3, 1.5, 7.4, 1.55,
                    fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY)
add_text(slide, 0.5, 1.6, 7.2, 1.35,
         "Palantir 帮助政府与大型企业将分散的数据、复杂的业务关系、规则、模型和人员行动连接在一起，"
         "在同一套系统中完成从「感知 → 决策 → 行动 → 结果反馈」的闭环，"
         "让组织在复杂和高风险的环境中更快、更准、更安全地行动。",
         size=11, color=COL_TEXT_GRAY, align=PP_ALIGN.LEFT)

# RIGHT: "为客户创造的核心价值" - 6 boxes in 2x3 grid
add_text(slide, 7.9, 1.4, 5.2, 0.35,
         "为客户创造的核心价值",
         size=13, color=COL_SECTION_HEADER, bold=True)

values = [
    ("打破数据孤岛", "形成单一态势"),
    ("更快更准的决策", "发现、模型 + AI"),
    ("让决策落地为行动", "产生真实影响"),
    ("企业级安全与治理", "可追溯、可审计"),
    ("持续学习与优化", "形成决策资产"),
    ("跨环境部署", "云·私有·边缘"),
]
for i, (title, sub) in enumerate(values):
    col = i % 3
    row = i // 3
    x = 7.9 + col * 1.78
    y = 1.8 + row * 0.7
    box = add_rect(slide, x, y, 1.7, 0.62,
                   fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
    add_text(slide, x + 0.05, y + 0.05, 1.6, 0.3,
             title, size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.05, y + 0.32, 1.6, 0.3,
             sub, size=9, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# ====================================================================
# ROW 2: Left "Palantir 的真实身份" 5 boxes + Right "Palantir 关键数据" 6 metrics
# ====================================================================
# LEFT section header
add_text(slide, 0.3, 3.25, 7.4, 0.35,
         "Palantir 的真实身份",
         size=14, color=COL_SECTION_HEADER, bold=True)

# 5 capability boxes
caps = [
    ("数据库", "数据连接与治理", "连接各数据源、清洗、可追溯"),
    ("业务建模", "Ontology 表达", "对象与关系刻画"),
    ("逻辑、智能", "决策能力建模", "规则 + 模型 + AI"),
    ("Action 行动", "Action 执行", "决策改变现实"),
    ("安全 + 权限", "安全与权限", "贯穿运营全过程"),
    ("反馈", "结果反馈", "持续学习与优化"),
]
for i, (icon, title, sub) in enumerate(caps):
    x = 0.3 + i * 1.23
    y = 3.7
    box = add_rect(slide, x, y, 1.13, 1.1,
                   fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=0.75)
    # Icon as small circle
    add_circle(slide, x + 0.42, y + 0.1, 0.3, 0.3,
               fill_color=COL_LIGHT_BLUE)
    add_text(slide, x + 0.42, y + 0.14, 0.3, 0.25,
             icon[:2] if len(icon) >= 2 else icon,
             size=10, color=COL_SECTION_HEADER, bold=True, align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, x + 0.02, y + 0.5, 1.09, 0.3,
             title, size=9, color=COL_TITLE_DARK, bold=True, align=PP_ALIGN.CENTER)
    # Subtitle
    add_text(slide, x + 0.02, y + 0.78, 1.09, 0.28,
             sub, size=7, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# RIGHT: Section header
add_text(slide, 7.9, 3.25, 5.2, 0.35,
         "Palantir 关键数据（截至 2026 年第一季度）",
         size=14, color=COL_SECTION_HEADER, bold=True)

# 6 metrics in 2 rows
metrics = [
    ("4,429", "全职员工", COL_ICON_BLUE),
    ("954+", "家客户", COL_GREEN),
    ("$44.75亿", "2025 年度收入", COL_RED_ACCENT),
    ("82%", "2025 毛利率", COL_SECTION_HEADER),
    ("85%", "2025 年毛利率", COL_MEDIUM_BLUE),
    ("87%", "2026 Q1 毛利率", COL_TITLE_DARK),
]
for i, (val, label, c) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = 7.9 + col * 1.78
    y = 3.7 + row * 0.7
    box = add_rect(slide, x, y, 1.7, 0.62,
                   fill_color=COL_VALUE_BG, line_color=COL_BORDER_GRAY, line_width=0.5)
    add_text(slide, x + 0.05, y + 0.05, 1.6, 0.32,
             val, size=16, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.05, y + 0.4, 1.6, 0.22,
             label, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)

# Source note
add_text(slide, 7.9, 5.18, 5.2, 0.25,
         "数据来源：Palantir 2025 年度报告、2026 年第一季度财报",
         size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.LEFT)

# ====================================================================
# ROW 3: Bottom flow "Palantir 连接什么 → 产生什么"
# ====================================================================
add_text(slide, 0.3, 5.55, 12.7, 0.35,
         "Palantir 连接什么 → 产生什么",
         size=14, color=COL_SECTION_HEADER, bold=True)

# 5 connected boxes
steps = [
    ("连接所有数据", "内部系统 (ERP, MES, CRM...)\n数据库 / 文档 / 传感器 / IoT\n外部数据 / 合作伙伴 / 公开数据", COL_SECTION_HEADER),
    ("表达业务\n现实世界", "客户\n工厂\n产品\n供应商", COL_ICON_BLUE),
    ("决策与智能", "业务规则\n预测模型\n优化算法", COL_GREEN),
    ("行动与执行", "审批与流程\n调用外部\n调整计划", COL_SECTION_HEADER),
    ("结果反馈与\n持续优化", "执行结果追踪\n效果分析\n沉淀资产", COL_RED_ACCENT),
]

box_w = 2.4
gap = 0.1
total_w = box_w * 5 + gap * 4
start_x = (13.33 - total_w) / 2

for i, (title, content, c) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    y = 6.0
    box = add_rect(slide, x, y, box_w, 1.05,
                   fill_color=COL_BG_WHITE, line_color=COL_BORDER_BLUE, line_width=1.0)
    add_text(slide, x + 0.05, y + 0.05, box_w - 0.1, 0.3,
             title, size=11, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.05, y + 0.32, box_w - 0.1, 0.7,
             content, size=8, color=COL_TEXT_GRAY, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        Inches(x + box_w - 0.05),
                                        Inches(y + 0.4),
                                        Inches(gap + 0.1),
                                        Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COL_ARROW_GRAY
        arrow.line.fill.background()
        arrow.shadow.inherit = False


# ====================================================================
# Bottom: ONE-LINE SUMMARY
# ====================================================================
summary_bg = add_rect(slide, 0.3, 7.15, 12.7, 0.3,
                       fill_color=COL_LIGHT_BLUE, line_color=COL_BORDER_BLUE, line_width=0.5)
add_text(slide, 0.4, 7.18, 12.5, 0.25,
         [("一句话总结：", 10, COL_TITLE_DARK, True),
          ("Palantir 不是传统软件厂商，而是「决策与运营基础设施」—— 通过统一的业务模型和动作体系，把数据、智能和行动连接起来，", 10, COL_TEXT_DARK, False),
          ("帮助客户在真实世界做准确决策并执行。", 10, COL_TITLE_DARK, True)])

# Save
prs.save(OUTPUT)
print(f"PPTX saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Width: {prs.slide_width/914400:.2f}\" Height: {prs.slide_height/914400:.2f}\"")
