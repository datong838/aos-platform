#!/usr/bin/env python3
"""
Generate PPTX from Toutiao article: "23页PPT讲透Palantir"
Embeds article images at corresponding sections.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

IMG_DIR = "/Users/ddt/WorkBuddy/2026-07-21-12-26-10/pptx_images"
OUTPUT = "/Users/ddt/work/projects/ai_agent/docs/ref/23页PPT讲透Palantir-Ontology-Data-Logic-Action-Security.pptx"

# Colors
DARK_BG = RGBColor(0x0F, 0x1A, 0x2E)       # Deep navy
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)   # Bright blue
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)   # Cyan
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xB0, 0xC0)
MID_GRAY = RGBColor(0x60, 0x70, 0x80)
GOLD = RGBColor(0xF0, 0xB9, 0x29)
GREEN = RGBColor(0x3A, 0xD9, 0x88)
RED_ACCENT = RGBColor(0xE8, 0x4D, 0x4D)
SECTION_BG = RGBColor(0x14, 0x21, 0x3D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG

def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox

def set_text(textbox, lines, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """lines: list of (text, size, color, bold) tuples"""
    tf = textbox.text_frame
    tf.clear()
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, size, clr, bld = item
        else:
            text, size, clr, bld = item, font_size, color, bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.name = "Microsoft YaHei"

def add_image_with_aspect(slide, img_path, left, top, max_width, max_height):
    """Add image preserving aspect ratio, centered in the bounding box"""
    if not os.path.exists(img_path):
        return None
    img = Image.open(img_path)
    w, h = img.size
    # Calculate aspect ratio
    aspect = w / h
    target_w = max_width
    target_h = int(target_w / aspect)
    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * aspect)
    # Center within bounding box
    actual_left = left + (max_width - target_w) // 2
    actual_top = top + (max_height - target_h) // 2
    return slide.shapes.add_picture(img_path, actual_left, actual_top, target_w, target_h)

def add_section_number(slide, num, total=17):
    """Add section number badge"""
    txBox = add_textbox(slide, Inches(0.5), Inches(0.2), Inches(2), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{num:02d} / {total:02d}"
    run.font.size = Pt(11)
    run.font.color.rgb = MID_GRAY
    run.font.name = "Microsoft YaHei"

def add_footer(slide, text="Palantir 深度解析"):
    txBox = add_textbox(slide, Inches(10), Inches(7.1), Inches(3), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MID_GRAY
    run.font.name = "Microsoft YaHei"


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Title
txBox = add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2))
set_text(txBox, [
    ("23页PPT讲透 Palantir", 42, WHITE, True),
], alignment=PP_ALIGN.LEFT)

# Subtitle
txBox = add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8))
set_text(txBox, [
    ("Ontology · Data · Logic · Action · Security", 24, ACCENT_CYAN, False),
], alignment=PP_ALIGN.LEFT)

# Description
txBox = add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(2))
set_text(txBox, [
    ("以 Ontology 为业务运行模型，以 Foundry 为数据与运营平台，", 16, LIGHT_GRAY, False),
    ("以 AIP 为 AI 运行环境，以 Apollo 和 Rubix 为软件交付底座，", 16, LIGHT_GRAY, False),
    ("以 FDE 为产品反馈机制，连接数据、决策和现实行动的", 16, LIGHT_GRAY, False),
    ("企业级决策与运营基础设施。", 18, GOLD, True),
], alignment=PP_ALIGN.LEFT)

# Date
txBox = add_textbox(slide, Inches(1.5), Inches(6.5), Inches(5), Inches(0.4))
p = txBox.text_frame.paragraphs[0]
run = p.add_run()
run.text = "2026.07 · 基于公开资料整理"
run.font.size = Pt(11)
run.font.color.rgb = MID_GRAY
run.font.name = "Microsoft YaHei"

# Embed cover image on right side
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_01.jpg"),
                      Inches(8.5), Inches(1.5), Inches(4), Inches(3))

# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 1)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.8))
set_text(txBox, [("目录 / Contents", 28, WHITE, True)])

contents_left = [
    "01  理解 Palantir：放弃简单标签",
    "02  发展历程：能力逐层叠加",
    "03  顶层产品体系：四大平台",
    "04  Foundry：数据运营平台",
    "05  Ontology：核心业务语言",
    "06  Action：从分析走向运营",
    "07  AIP：企业 AI 运行环境",
    "08  Apollo & Rubix：交付底座",
]
contents_right = [
    "09  Global Branching：软件工程延伸",
    "10  Gotham：任务运营系统",
    "11  FDE：产品研发体系",
    "12  Bootcamp：销售与验证",
    "13  商业模式：非纯 SaaS",
    "14  经营结果：增长数据",
    "15  真正壁垒：复合飞轮",
    "16  约束与挑战",
    "17  最终结论",
]

txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
set_text(txBox, [(c, 14, LIGHT_GRAY, False) for c in contents_left])

txBox = add_textbox(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5))
set_text(txBox, [(c, 14, LIGHT_GRAY, False) for c in contents_right])

add_footer(slide)

# ============================================================
# SLIDE 3: §1 理解 Palantir
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 2)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("01  理解 Palantir：首先要放弃几个简单标签", 26, ACCENT_CYAN, True)])

txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(6), Inches(1))
set_text(txBox, [
    ("常见但不够准确的标签：", 14, LIGHT_GRAY, False),
    ("大数据分析 · 商业智能 · 数据中台 · AI软件公司", 13, MID_GRAY, False),
    ("政府承包商 · 咨询公司 · Ontology 平台", 13, MID_GRAY, False),
])

txBox = add_textbox(slide, Inches(0.8), Inches(3.5), Inches(7), Inches(2.5))
set_text(txBox, [
    ("更准确的定义：", 14, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("Palantir 是一家以数据和软件工程为基础，", 20, GOLD, True),
    ("为政府及大型企业建设实时决策与运营", 20, GOLD, True),
    ("系统的企业软件公司。", 20, GOLD, True),
])

# Right side: seven dimensions
txBox = add_textbox(slide, Inches(8), Inches(1.8), Inches(4.8), Inches(5))
set_text(txBox, [
    ("研究 Palantir 必须同时看 7 个维度：", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("1. 如何连接和治理数据", 13, LIGHT_GRAY, False),
    ("2. 如何表达企业业务", 13, LIGHT_GRAY, False),
    ("3. 如何把规则引入决策", 13, LIGHT_GRAY, False),
    ("4. 如何把决策转化为行动", 13, LIGHT_GRAY, False),
    ("5. 如何在复杂环境中部署", 13, LIGHT_GRAY, False),
    ("6. 如何进入客户并扩张", 13, LIGHT_GRAY, False),
    ("7. 如何将经验沉淀为产品", 13, LIGHT_GRAY, False),
])

add_footer(slide)

# ============================================================
# SLIDE 4: §2 发展历程
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 3)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("02  发展不是转型，而是能力逐层叠加", 26, ACCENT_CYAN, True)])

# Four phases as timeline
phases = [
    ("2003", "Phase 1", "Gotham", "反恐情报 + 复杂关系分析\n共同态势认知 + 行动"),
    ("2010", "Phase 2", "Foundry", "商业企业运营\n50+ 垂直领域"),
    ("2018", "Phase 3", "Apollo", "多云/私有化/边缘部署\n软件持续交付"),
    ("2023", "Phase 4", "AIP", "生成式 AI 接入业务运营\nLLM + Agent + 自动化"),
]

x_start = 0.8
for i, (year, phase, product, desc) in enumerate(phases):
    x = Inches(x_start + i * 3.1)

    # Year
    txBox = add_textbox(slide, x, Inches(1.8), Inches(2.8), Inches(0.5))
    set_text(txBox, [(year, 28, GOLD, True)])

    # Phase badge
    txBox = add_textbox(slide, x, Inches(2.4), Inches(2.8), Inches(0.4))
    set_text(txBox, [(phase, 11, MID_GRAY, False)])

    # Product name
    txBox = add_textbox(slide, x, Inches(2.9), Inches(2.8), Inches(0.5))
    set_text(txBox, [(product, 22, ACCENT_CYAN, True)])

    # Arrow
    txBox = add_textbox(slide, x, Inches(3.5), Inches(2.8), Inches(0.4))
    set_text(txBox, [("→", 20, MID_GRAY, False)])

    # Description
    txBox = add_textbox(slide, x, Inches(4.0), Inches(2.8), Inches(2))
    set_text(txBox, [(desc, 12, LIGHT_GRAY, False)])

# Bottom insight
txBox = add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8))
set_text(txBox, [
    ("核心洞察：不是「数据分析公司转型 AI 公司」，而是在原有数据/业务模型/安全/运营之上增加 AI 这一新的决策主体。", 13, ACCENT_BLUE, True),
])

add_footer(slide)

# ============================================================
# SLIDE 5: §3 顶层产品体系 (with img_02)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 4)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("03  顶层产品体系：四大平台", 26, ACCENT_CYAN, True)])

# Product table
products = [
    ("Gotham", "政府/国防/情报", "情报分析、军事行动、态势感知", RED_ACCENT),
    ("Foundry", "数据运营 + 企业应用", "数据集成、业务建模、分析、工作流", ACCENT_BLUE),
    ("Apollo", "软件持续交付", "多云、私有化、边缘和隔离环境部署", GREEN),
    ("AIP", "企业生成式 AI", "LLM、Agent、自动化、评测、AI 应用", GOLD),
]

for i, (name, role, apps, color) in enumerate(products):
    y = Inches(1.8 + i * 1.15)
    # Name
    txBox = add_textbox(slide, Inches(0.8), y, Inches(2), Inches(0.5))
    set_text(txBox, [(name, 18, color, True)])
    # Role
    txBox = add_textbox(slide, Inches(3.0), y, Inches(2.5), Inches(0.5))
    set_text(txBox, [(role, 13, WHITE, False)])
    # Apps
    txBox = add_textbox(slide, Inches(5.8), y, Inches(4), Inches(0.8))
    set_text(txBox, [(apps, 12, LIGHT_GRAY, False)])

# Key note about Ontology
txBox = add_textbox(slide, Inches(0.8), Inches(6.2), Inches(7), Inches(0.8))
set_text(txBox, [
    ("⚠ Ontology 不是第五个平台，而是贯穿 Foundry 和 AIP 的核心架构层", 13, GOLD, True),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_02.jpg"),
                      Inches(9.5), Inches(1.5), Inches(3.2), Inches(4))

add_footer(slide)

# ============================================================
# SLIDE 6: §4 Foundry
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 5)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(9), Inches(0.8))
set_text(txBox, [("04  Foundry：不是数据仓库，而是数据运营平台", 26, ACCENT_CYAN, True)])

# Three key points
txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.5))
set_text(txBox, [
    ("职责", 16, ACCENT_BLUE, True),
    ("", 6, WHITE, False),
    ("数据连接", 13, WHITE, False),
    ("数据转换和工程", 13, LIGHT_GRAY, False),
    ("数据质量 + 血缘", 13, LIGHT_GRAY, False),
    ("数据权限", 13, LIGHT_GRAY, False),
    ("业务逻辑开发", 13, LIGHT_GRAY, False),
    ("Ontology 建模", 13, LIGHT_GRAY, False),
    ("模型运行", 13, LIGHT_GRAY, False),
    ("应用开发 + 工作流", 13, LIGHT_GRAY, False),
    ("自动化", 13, LIGHT_GRAY, False),
])

txBox = add_textbox(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.5))
set_text(txBox, [
    ("Multimodal Data Plane", 16, ACCENT_BLUE, True),
    ("", 6, WHITE, False),
    ("结构化数据 + 非结构化", 13, WHITE, False),
    ("流数据 + 地理空间", 13, LIGHT_GRAY, False),
    ("文档 + 图片 + 视频", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("连接 ERP / CRM / 数据仓库", 13, LIGHT_GRAY, False),
    ("工业数据库 / 传感器", 13, LIGHT_GRAY, False),
    ("与 Snowflake / Databricks 互操作", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("不要求所有数据进入一个仓库", 13, GOLD, True),
])

txBox = add_textbox(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
set_text(txBox, [
    ("运营统一 ≠ 数据统一", 16, ACCENT_BLUE, True),
    ("", 6, WHITE, False),
    ("传统终点：", 13, MID_GRAY, False),
    ("  数据集 / 指标 / 报表", 13, LIGHT_GRAY, False),
    ("  模型输出", 13, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("Foundry 终点：", 13, WHITE, False),
    ("  业务对象 / 运营应用", 13, LIGHT_GRAY, False),
    ("  任务 / 工作流 / Action", 13, LIGHT_GRAY, False),
    ("  业务系统写回", 13, LIGHT_GRAY, False),
])

add_footer(slide)

# ============================================================
# SLIDE 7: §5 Ontology 核心业务语言 (with img_03)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 6)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("05  Ontology：核心业务语言", 26, ACCENT_CYAN, True)])

# Key quote
txBox = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(1))
set_text(txBox, [
    ("\"Ontology 被设计用于表达企业复杂而相互关联的决策，", 15, GOLD, True),
    (" 而不仅仅是表达数据。\"", 15, GOLD, True),
])

# Four core components
components = [
    ("Data", "决策依赖的事实", "当前状态 · 历史 · 实时 · 外部", ACCENT_BLUE),
    ("Logic", "计算和推理", "规则 · 指标 · ML模型 · LLM函数", GREEN),
    ("Action", "影响现实世界", "修改对象 · 审批 · 写回系统", GOLD),
    ("Security", "权限控制", "查看 · 运行 · 发起 · 修改", RED_ACCENT),
]

for i, (name, role, detail, color) in enumerate(components):
    x = Inches(0.8 + i * 3.1)
    # Component name
    txBox = add_textbox(slide, x, Inches(3.0), Inches(2.8), Inches(0.5))
    set_text(txBox, [(name, 20, color, True)])
    # Role
    txBox = add_textbox(slide, x, Inches(3.6), Inches(2.8), Inches(0.4))
    set_text(txBox, [(role, 13, WHITE, False)])
    # Detail
    txBox = add_textbox(slide, x, Inches(4.1), Inches(2.8), Inches(0.8))
    set_text(txBox, [(detail, 11, LIGHT_GRAY, False)])

# Noun & Verb metaphor
txBox = add_textbox(slide, Inches(0.8), Inches(5.2), Inches(7), Inches(1.5))
set_text(txBox, [
    ("名词（对象 + 关系）→ 描述企业", 14, LIGHT_GRAY, False),
    ("动词（Action）→ 改变企业", 14, GOLD, True),
    ("", 6, WHITE, False),
    ("Ontology 不是静态知识图谱，而是企业业务对象运行时。", 13, ACCENT_CYAN, False),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_03.jpg"),
                      Inches(8.8), Inches(3.5), Inches(3.8), Inches(3.5))

add_footer(slide)

# ============================================================
# SLIDE 8: §6 Action 从分析走向运营
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 7)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(9), Inches(0.8))
set_text(txBox, [("06  Action：从分析走向运营的关键", 26, ACCENT_CYAN, True)])

# Traditional vs Palantir flow
txBox = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5))
set_text(txBox, [
    ("传统数据/AI 平台流程：", 14, MID_GRAY, True),
    ("发现问题 → 提供建议 → 用户离开平台执行", 13, LIGHT_GRAY, False),
])

txBox = add_textbox(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.5))
set_text(txBox, [
    ("Palantir 流程：", 14, GOLD, True),
    ("发现问题 → 生成方案 → 权限校验", 13, WHITE, False),
    ("→ 执行 Action → 结果反馈", 13, WHITE, False),
])

# Action capabilities
txBox = add_textbox(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4))
set_text(txBox, [
    ("Action 包含：", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("• 输入参数 + 表单", 13, LIGHT_GRAY, False),
    ("• 业务规则 + 提交条件", 13, LIGHT_GRAY, False),
    ("• 权限控制", 13, LIGHT_GRAY, False),
    ("• 对象修改（一个事务中修改多个对象）", 13, LIGHT_GRAY, False),
    ("• 外部副作用（写回业务系统）", 13, LIGHT_GRAY, False),
    ("• 审计 + 运行指标", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("语义 + 动力学：", 14, GOLD, True),
    ("语义说明企业是什么；", 13, LIGHT_GRAY, False),
    ("动力学说明企业如何发生变化。", 13, LIGHT_GRAY, False),
])

# Feedback loop
txBox = add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.5))
set_text(txBox, [
    ("决策反馈循环：", 14, ACCENT_BLUE, True),
    ("Action 结果 → 重新进入 Ontology → 调整规则 / 重训模型 / 优化建议", 12, LIGHT_GRAY, False),
])

add_footer(slide)

# ============================================================
# SLIDE 9: §7 AIP 企业AI运行环境 (with img_04)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 8)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("07  AIP：不是大模型网关，而是企业 AI 运行环境", 24, ACCENT_CYAN, True)])

# Without vs With Ontology
txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5))
set_text(txBox, [
    ("没有 Foundry + Ontology：", 14, RED_ACCENT, True),
    ("", 6, WHITE, False),
    ("大模型只能看到：", 13, LIGHT_GRAY, False),
    ("• 数据库字段", 12, MID_GRAY, False),
    ("• 搜索结果", 12, MID_GRAY, False),
    ("• 文档片段", 12, MID_GRAY, False),
    ("• Prompt 临时上下文", 12, MID_GRAY, False),
])

txBox = add_textbox(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5))
set_text(txBox, [
    ("接入 Ontology 后：", 14, GREEN, True),
    ("", 6, WHITE, False),
    ("AI 可以知道：", 13, WHITE, False),
    ("• 当前对象是什么 + 关系", 12, LIGHT_GRAY, False),
    ("• 可以调用哪些规则和模型", 12, LIGHT_GRAY, False),
    ("• 可以执行哪些 Action", 12, LIGHT_GRAY, False),
    ("• 当前用户拥有什么权限", 12, LIGHT_GRAY, False),
])

# AIP capabilities
txBox = add_textbox(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4))
set_text(txBox, [
    ("AIP 核心能力：", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("• 多模型连接 + Agent 开发", 12, LIGHT_GRAY, False),
    ("• 自动化 + AI 应用", 12, LIGHT_GRAY, False),
    ("• 上下文工程 + 工具调用", 12, LIGHT_GRAY, False),
    ("• AI 评测 + 可观测", 12, LIGHT_GRAY, False),
    ("• 安全治理 + 人工审核", 12, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("Agent 继承现实权限：", 14, GOLD, True),
    ("用户不能创建代码库，AI 也不能；", 12, LIGHT_GRAY, False),
    ("用户不能修改对象，AI 也不能。", 12, LIGHT_GRAY, False),
    ("AI 不绕过权限体系，", 12, ACCENT_CYAN, False),
    ("而是权限体系中的受控参与者。", 12, ACCENT_CYAN, False),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_04.jpg"),
                      Inches(10.5), Inches(1.5), Inches(2.2), Inches(2.5))

add_footer(slide)

# ============================================================
# SLIDE 10: §8 Apollo & Rubix (with img_05)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 9)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("08  Apollo & Rubix：经常被低估的壁垒", 26, ACCENT_CYAN, True)])

# Apollo
txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4))
set_text(txBox, [
    ("Apollo — 持续交付", 18, ACCENT_BLUE, True),
    ("", 6, WHITE, False),
    ("每天协调数千次跨数百项服务", 13, WHITE, False),
    ("的零停机升级", 13, WHITE, False),
    ("", 6, WHITE, False),
    ("• 软件版本管理", 12, LIGHT_GRAY, False),
    ("• 配置管理 + 多环境发布", 12, LIGHT_GRAY, False),
    ("• 安全更新 + 灰度部署", 12, LIGHT_GRAY, False),
    ("• 回滚 + 持续运维", 12, LIGHT_GRAY, False),
])

# Rubix
txBox = add_textbox(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4))
set_text(txBox, [
    ("Rubix — 运行底座", 18, GREEN, True),
    ("", 6, WHITE, False),
    ("强化后的 Kubernetes 运行底座", 13, WHITE, False),
    ("", 6, WHITE, False),
    ("• 高可用 + 自动扩缩容", 12, LIGHT_GRAY, False),
    ("• 工作负载隔离 + 默认安全网络", 12, LIGHT_GRAY, False),
    ("• 加密 + 身份验证 + 授权", 12, LIGHT_GRAY, False),
    ("• 临时计算节点 + 多环境一致", 12, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("支持部署到 AWS / Azure / GCP / Oracle", 12, LIGHT_GRAY, False),
    ("/ 私有数据中心 / 高安全隔离环境", 12, LIGHT_GRAY, False),
])

# Bottom insight
txBox = add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1))
set_text(txBox, [
    ("为什么这一层重要：企业客户需要软件长期运行、持续升级、满足安全要求、适配已有基础设施——", 13, GOLD, False),
    ("Apollo 和 Rubix 解决的是「从可以开发」到「可以长期运营」的跨越。", 13, ACCENT_CYAN, True),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_05.jpg"),
                      Inches(9.5), Inches(5.2), Inches(3), Inches(2))

add_footer(slide)

# ============================================================
# SLIDE 11: §9 Global Branching (with img_06)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 10)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("09  Global Branching：将软件工程延伸到整个业务系统", 24, ACCENT_CYAN, True)])

# Traditional Git vs Palantir Global Branching
txBox = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5), Inches(1.5))
set_text(txBox, [
    ("传统 Git 管理：", 14, MID_GRAY, True),
    ("", 6, WHITE, False),
    ("仅管理代码", 13, LIGHT_GRAY, False),
])

txBox = add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(3))
set_text(txBox, [
    ("Global Branching 管理：", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("• 数据管道", 13, WHITE, False),
    ("• 数据 Schema", 13, WHITE, False),
    ("• Ontology（对象/属性/关系）", 13, WHITE, False),
    ("• Action", 13, WHITE, False),
    ("• 应用 + Functions", 13, WHITE, False),
    ("• 自动化流程", 13, WHITE, False),
])

txBox = add_textbox(slide, Inches(6.5), Inches(2.5), Inches(6), Inches(4))
set_text(txBox, [
    ("隔离分支开发流程：", 14, GOLD, True),
    ("", 6, WHITE, False),
    ("创建分支 → 修改完整工作流", 13, LIGHT_GRAY, False),
    ("→ 测试（Action 可运行但不写主分支）", 13, LIGHT_GRAY, False),
    ("→ 审核通过 → 统一合并", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("企业不仅需要管理代码变化，", 13, ACCENT_BLUE, False),
    ("还需要管理业务定义和决策流程的变化。", 13, ACCENT_CYAN, True),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_06.jpg"),
                      Inches(10), Inches(5.5), Inches(2.5), Inches(1.8))

add_footer(slide)

# ============================================================
# SLIDE 12: §10-12 Gotham / FDE / Bootcamp (with img_07)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 11)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("10-12  Gotham · FDE · Bootcamp", 26, ACCENT_CYAN, True)])

# Three columns
# Gotham
txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5))
set_text(txBox, [
    ("Gotham", 18, RED_ACCENT, True),
    ("任务运营系统", 13, WHITE, False),
    ("", 6, WHITE, False),
    ("• 情报 + 传感器数据整合", 12, LIGHT_GRAY, False),
    ("• 地理空间 + 目标 + 人员", 12, LIGHT_GRAY, False),
    ("• 任务规划 + 资产调度", 12, LIGHT_GRAY, False),
    ("• 传感器控制（无人机/卫星）", 12, LIGHT_GRAY, False),
    ("• 人在回路控制", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("国防任务操作系统，", 12, GOLD, False),
    ("而非只做情报分析", 12, GOLD, False),
])

# FDE
txBox = add_textbox(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5))
set_text(txBox, [
    ("FDE", 18, ACCENT_BLUE, True),
    ("前向部署工程", 13, WHITE, False),
    ("", 6, WHITE, False),
    ("\"人类版本的反向传播\"", 12, GOLD, True),
    ("", 6, WHITE, False),
    ("• 业务理解 + 数据建模", 12, LIGHT_GRAY, False),
    ("• 软件开发 + 产品配置", 12, LIGHT_GRAY, False),
    ("• 系统集成 + 场景验证", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("不是实施工程师，", 12, ACCENT_CYAN, False),
    ("而是产品研发的一部分", 12, ACCENT_CYAN, True),
])

# Bootcamp
txBox = add_textbox(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5))
set_text(txBox, [
    ("Bootcamp", 18, GREEN, True),
    ("验证 + 交付 + 销售", 13, WHITE, False),
    ("", 6, WHITE, False),
    ("• 使用客户真实数据", 12, LIGHT_GRAY, False),
    ("• 客户人员直接参与", 12, LIGHT_GRAY, False),
    ("• 数小时/数天形成可运行工作流", 12, LIGHT_GRAY, False),
    ("• 不做通用产品演示", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("2024年超 500 场", 12, GOLD, False),
    ("同时承担：", 12, ACCENT_CYAN, False),
    ("产品验证 / 用例发现 /", 11, LIGHT_GRAY, False),
    ("客户培训 / 扩张入口", 11, LIGHT_GRAY, False),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_07.jpg"),
                      Inches(5), Inches(6.5), Inches(3.5), Inches(0.8))

add_footer(slide)

# ============================================================
# SLIDE 13: §13 商业模式 (with img_08)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 12)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("13  商业模式：不是纯 SaaS", 26, ACCENT_CYAN, True)])

# Revenue composition
txBox = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5))
set_text(txBox, [
    ("收入构成：", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("• Palantir Cloud 订阅", 12, LIGHT_GRAY, False),
    ("• 本地/私有部署软件订阅", 12, LIGHT_GRAY, False),
    ("• 持续运营和维护服务", 12, LIGHT_GRAY, False),
    ("• 专业服务（培训/配置/建模）", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("= 平台订阅 + 持续运维 + 高价值工程服务", 12, GOLD, True),
])

# Land and Expand
txBox = add_textbox(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.5))
set_text(txBox, [
    ("Land and Expand：", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("一个用例 → 一个部门", 12, LIGHT_GRAY, False),
    ("→ 更多数据源 → 更多对象", 12, LIGHT_GRAY, False),
    ("→ 更多用户 → 更多工作流", 12, LIGHT_GRAY, False),
    ("→ 更大平台合同", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("2025年 Top20 客户均收入 9,390 万美元", 12, GOLD, True),
    ("同比增长 45%", 12, GOLD, True),
])

# Government vs Commercial
txBox = add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(2.5))
set_text(txBox, [
    ("政府 + 商业双轮（2025年 44.75 亿美元）", 14, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("政府收入 24.02 亿美元（54%）← 高安全 + 任务关键 ←→ 商业收入 20.73 亿美元（46%）← 行业场景 + 增长空间", 13, WHITE, False),
    ("共用平台让两个市场的技术积累相互复用", 13, LIGHT_GRAY, False),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_08.jpg"),
                      Inches(9.5), Inches(5), Inches(3), Inches(2))

add_footer(slide)

# ============================================================
# SLIDE 14: §14 经营结果 (with img_09)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 13)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("14  当前经营结果", 26, ACCENT_CYAN, True)])

# Key metrics
metrics = [
    ("4,429", "全职员工", ACCENT_CYAN),
    ("954+", "客户数量", ACCENT_BLUE),
    ("$44.75亿", "2025年度收入", GOLD),
    ("82%", "毛利率", GREEN),
    ("$14.14亿", "营业利润", GOLD),
    ("85%", "Q1'26 收入增长", GREEN),
]

for i, (val, label, color) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.0 + row * 2.0)

    txBox = add_textbox(slide, x, y, Inches(3.5), Inches(0.8))
    set_text(txBox, [(val, 32, color, True)])

    txBox = add_textbox(slide, x, y + Inches(0.8), Inches(3.5), Inches(0.5))
    set_text(txBox, [(label, 14, LIGHT_GRAY, False)])

# Key insights
txBox = add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8))
set_text(txBox, [
    ("结论：政府+商业并重 · AIP 加速商业增长 · 高软件毛利 · 既有客户扩张是重要来源", 13, ACCENT_CYAN, True),
])

# Embed image
add_image_with_aspect(slide, os.path.join(IMG_DIR, "img_09.jpg"),
                      Inches(9.5), Inches(6), Inches(3), Inches(1.5))

add_footer(slide)

# ============================================================
# SLIDE 15: §15 真正壁垒 — 复合飞轮
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 14)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("15  真正壁垒：一套相互增强的系统", 26, ACCENT_CYAN, True)])

# Flywheel components
flywheel = [
    ("数据 + 计算", "连接复杂异构高敏感数据"),
    ("Ontology", "统一业务运行模型"),
    ("Action + 反馈循环", "分析结果进入业务执行"),
    ("AIP", "大模型进入已有业务模型"),
    ("Apollo + Rubix", "多云/边缘部署 + 持续升级"),
    ("FDE", "客户问题转化为产品能力"),
    ("Bootcamp", "降低初次验证成本"),
    ("账户扩张", "一次性问题 → 长期平台收入"),
]

for i, (name, desc) in enumerate(flywheel):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.8 + row * 0.7)

    txBox = add_textbox(slide, x, y, Inches(2.5), Inches(0.5))
    set_text(txBox, [(name, 14, GOLD, True)])

    txBox = add_textbox(slide, x + Inches(2.8), y, Inches(3), Inches(0.5))
    set_text(txBox, [(desc, 12, LIGHT_GRAY, False)])

# Flywheel equation
txBox = add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(2))
set_text(txBox, [
    ("复合飞轮：", 14, ACCENT_CYAN, True),
    ("", 4, WHITE, False),
    ("更多复杂客户问题 → 更多现场工程经验 → 更强的平台能力 → 更快的用例交付", 12, LIGHT_GRAY, False),
    ("→ 更多客户和场景 → 更多 Ontology/Action/决策数据 → 更高平台价值 → 更深客户扩张", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("Palantir 真正的壁垒不是一项技术专利，而是产品·工程·交付·客户·业务数据共同形成的复合飞轮。", 14, GOLD, True),
])

add_footer(slide)

# ============================================================
# SLIDE 16: §16 约束与挑战
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 15)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("16  也存在明显约束", 26, ACCENT_CYAN, True)])

constraints = [
    ("实施复杂度高", "平台复杂，实施周期长，需要针对客户独特环境配置 + 培训 + 持续技术人员服务"),
    ("销售成本高", "大型政府/商业项目安装成本高、失败风险高、评估周期可能超过一年"),
    ("依赖工程服务", "成本含现场人员/专业服务/分包商/云资源/部署维护。部分客户仍需长期参与"),
    ("面临多类竞争", "大型软件公司 / 政府承包商 / 系统集成商 / 新兴技术公司 / 客户自研系统"),
]

for i, (title, desc) in enumerate(constraints):
    y = Inches(1.8 + i * 1.3)
    txBox = add_textbox(slide, Inches(0.8), y, Inches(3), Inches(0.5))
    set_text(txBox, [(title, 15, RED_ACCENT, True)])

    txBox = add_textbox(slide, Inches(0.8), y + Inches(0.5), Inches(11), Inches(0.7))
    set_text(txBox, [(desc, 13, LIGHT_GRAY, False)])

txBox = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5))
set_text(txBox, [
    ("竞争优势不是某个模块全面领先，而是能把多个模块组合成一套完整运营系统。", 13, GOLD, False),
])

add_footer(slide)

# ============================================================
# SLIDE 17: §17 最终结论
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_section_number(slide, 16)

txBox = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
set_text(txBox, [("17  最终研究结论", 26, ACCENT_CYAN, True)])

conclusions = [
    ("一", "Palantir 不是传统数据公司", "数据是起点不是终点 — 最终处理的是数据→决策→行动→改变现实"),
    ("二", "Ontology 是核心但不是全部", "只有与数据/逻辑/Action/AI/安全/部署结合后才具运营价值"),
    ("三", "AIP 不是替代而是放大器", "用大模型提高企业系统的自然语言交互/推理/工具调用/自动化"),
    ("四", "Apollo 决定能否进入关键环境", "能部署到大型企业/政府/军事/私有环境/边缘现场并持续维护"),
    ("五", "FDE 是产品研发方法", "客户现场是产品研发循环的一部分"),
    ("六", "商业模式建立在账户扩张", "从复杂客户高投入开始，逐渐进入更多部门/数据/运营流程"),
]

for i, (num, title, desc) in enumerate(conclusions):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.8 + row * 1.5)

    txBox = add_textbox(slide, x, y, Inches(0.5), Inches(0.5))
    set_text(txBox, [(num, 22, GOLD, True)])

    txBox = add_textbox(slide, x + Inches(0.6), y, Inches(5), Inches(0.5))
    set_text(txBox, [(title, 14, WHITE, True)])

    txBox = add_textbox(slide, x + Inches(0.6), y + Inches(0.5), Inches(5), Inches(0.8))
    set_text(txBox, [(desc, 11, LIGHT_GRAY, False)])

# Final definition
txBox = add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8))
set_text(txBox, [
    ("结论七：Palantir = 企业级决策与运营基础设施", 16, GOLD, True),
], alignment=PP_ALIGN.CENTER)

add_footer(slide)

# ============================================================
# SLIDE 18: CLOSING — One-line definition
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Center text
txBox = add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "以 Ontology 为业务运行模型"
run.font.size = Pt(20)
run.font.color.rgb = LIGHT_GRAY
run.font.name = "Microsoft YaHei"

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "以 Foundry 为数据与运营平台"
run2.font.size = Pt(20)
run2.font.color.rgb = LIGHT_GRAY
run2.font.name = "Microsoft YaHei"

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "以 AIP 为 AI 运行环境"
run3.font.size = Pt(20)
run3.font.color.rgb = LIGHT_GRAY
run3.font.name = "Microsoft YaHei"

p4 = tf.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "以 Apollo 和 Rubix 为软件交付底座"
run4.font.size = Pt(20)
run4.font.color.rgb = LIGHT_GRAY
run4.font.name = "Microsoft YaHei"

p5 = tf.add_paragraph()
p5.space_before = Pt(20)
p5.alignment = PP_ALIGN.CENTER
run5 = p5.add_run()
run5.text = "连接数据、决策和现实行动的"
run5.font.size = Pt(18)
run5.font.color.rgb = ACCENT_CYAN
run5.font.name = "Microsoft YaHei"

p6 = tf.add_paragraph()
p6.alignment = PP_ALIGN.CENTER
run6 = p6.add_run()
run6.text = "企业级决策与运营基础设施"
run6.font.size = Pt(28)
run6.font.color.rgb = GOLD
run6.font.bold = True
run6.font.name = "Microsoft YaHei"

# Source
txBox = add_textbox(slide, Inches(3), Inches(6.5), Inches(7), Inches(0.5))
p = txBox.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "来源：今日头条「23页PPT讲透Palantir」 · 2026.07"
run.font.size = Pt(10)
run.font.color.rgb = MID_GRAY
run.font.name = "Microsoft YaHei"


# Save
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"PPTX saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
print(f"Images embedded: 9")
